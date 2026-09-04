import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors matching the sleek, modern PM case study design
NAVY = RGBColor(15, 23, 42)        # #0F172A Dark Slate Header Text
ACCENT_BLUE = RGBColor(27, 38, 79) # #1B264F Dark Navy Container
CORAL = RGBColor(244, 63, 94)      # #F43F5E Pinkish Coral Header Accent
CORAL_DARK = RGBColor(225, 29, 72) # #E11D48 Accent Red
GOLD = RGBColor(217, 119, 6)      # #D97706 Warm Gold Accent
LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC Card BG
CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0 Slate border
TEXT_MAIN = RGBColor(15, 23, 42)  # #0F172A Primary Text
TEXT_MUTED = RGBColor(71, 85, 105) # #475569 Muted Text
WHITE = RGBColor(255, 255, 255)
LIGHT_RED_BG = RGBColor(255, 241, 242) # #FFF1F2
LIGHT_GREEN_BG = RGBColor(240, 253, 244) # #F0FDF4
LIGHT_BLUE_BG = RGBColor(239, 246, 255) # #EFF6FF

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen (13.333 inches x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_slide_header(slide, category, title, subtitle=""):
        # Background canvas
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()

        # Category Kicker (14pt bold coral)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.35))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = CORAL_DARK

        # Actionable Takeaway Slide Title (24pt bold)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.75))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Subtitle if present (14pt muted)
        if subtitle:
            tx_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.4))
            tf_sub = tx_sub.text_frame
            tf_sub.word_wrap = True
            p_s = tf_sub.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(14)
            p_s.font.color.rgb = TEXT_MUTED

        # Anonymous Footer (14pt)
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_f = footer.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "Myntra · Wishlist-to-Purchase Growth Case Study"
        p_f.font.size = Pt(14)
        p_f.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: Hero & Strategic Thesis
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = WHITE
    bg1.line.fill.background()

    # Kicker
    tx = slide1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(7.5), Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "GROWTH PM CASE STUDY · MYNTRA"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CORAL_DARK

    # Title
    tx = slide1.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(7.5), Inches(1.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Convert more wishlisters by resolving the last uncertainty — not by discounting."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Subtitle
    tx = slide1.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(7.5), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Smart Wishlist + FitCheck turns a passive save into a decision-ready state for high-intent apparel shoppers — and measures whether that confidence actually moves them to Bag and purchase within 30 days."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED

    # 3 Stat Cards
    stat_data = [
        ("51.3%", "of analyzed conversations fell in high-intent / high-friction", CORAL_DARK),
        ("23/25", "Fit uncertainty ranked #1 on the opportunity rubric", GOLD),
        ("5 / 6", "interviews surfaced fit / drape uncertainty", RGBColor(16, 185, 129))
    ]
    for i, (num, desc, bar_col) in enumerate(stat_data):
        l = Inches(0.8) + i * Inches(2.55)
        c = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, Inches(3.9), Inches(2.4), Inches(1.4))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_BG
        c.line.color.rgb = CARD_BORDER

        # Left Accent Line
        line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(3.9), Inches(0.1), Inches(1.4))
        line.fill.solid()
        line.fill.fore_color.rgb = bar_col
        line.line.fill.background()

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

    # Buttons
    btn1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(3.6), Inches(0.65))
    btn1.fill.solid()
    btn1.fill.fore_color.rgb = ACCENT_BLUE
    btn1.line.fill.background()
    tf_b1 = btn1.text_frame
    p = tf_b1.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Open AI Discovery Engine 🔗"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    btn2 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(5.5), Inches(3.6), Inches(0.65))
    btn2.fill.solid()
    btn2.fill.fore_color.rgb = CORAL_DARK
    btn2.line.fill.background()
    tf_b2 = btn2.text_frame
    p = tf_b2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Open Live MVP 🚀"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Footnote
    tx = slide1.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(7.5), Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "No monetary incentives · 10-slide submission · live artefacts linked"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED

    # Right Container (THE MECHANISM PANEL)
    mech = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(0.8), Inches(4.0), Inches(5.9))
    mech.fill.solid()
    mech.fill.fore_color.rgb = LIGHT_BG
    mech.line.color.rgb = CARD_BORDER

    tf_m = mech.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = Inches(0.25)
    tf_m.margin_top = Inches(0.25)

    p = tf_m.paragraphs[0]
    p.text = "THE MECHANISM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CORAL_DARK

    p2 = tf_m.add_paragraph()
    p2.text = "From “Saved” to “Ready to Buy”"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = NAVY

    steps = [
        ("1", "SAVED", "User has explicit interest, not commitment"),
        ("2", "FIT UNCERTAIN", "Cross-brand size / drape is unresolved"),
        ("3", "FIT RESOLVED", "FitCheck returns transparent signal"),
        ("4", "READY TO BUY", "User decides; item becomes action-ready"),
        ("5", "BAG → PURCHASE", "Commitment advances toward 30-day goal")
    ]
    for i, (num, label, desc) in enumerate(steps):
        top_y = Inches(1.8) + i * Inches(0.85)
        # Circle badge
        circ = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), top_y, Inches(0.4), Inches(0.4))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE if i != 2 and i != 3 else CORAL_DARK
        circ.line.fill.background()
        tf_c = circ.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE

        # Text
        tx = slide1.shapes.add_textbox(Inches(9.3), top_y - Inches(0.1), Inches(3.1), Inches(0.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

    # Bottom Thesis Banner inside right panel
    th_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(6.0), Inches(3.6), Inches(0.6))
    th_box.fill.solid()
    th_box.fill.fore_color.rgb = LIGHT_RED_BG
    th_box.line.color.rgb = CORAL_DARK
    tf_th = th_box.text_frame
    tf_th.word_wrap = True
    p = tf_th.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thesis: decision confidence is the addressable growth lever."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CORAL_DARK

    # =========================================================================
    # SLIDE 2: Business Metric Decomposition
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "Business Metric Decomposition", "The business metric moves only when wishlisted users progress from intent to commitment.", "The brief is user-level, so the North Star must also be user-level — not an SKU conversion proxy.")

    # Formula Box
    f_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.733), Inches(0.75))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = ACCENT_BLUE
    f_box.line.fill.background()
    tf = f_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "30-Day Wishlist Buyer Conversion  =  (Unique users who purchase ≥1 item they wishlisted within 30 days of that add) ÷ (Unique users with ≥1 eligible wishlist add in cohort window)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Levers header badge
    l_badge = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(1.2), Inches(0.35))
    l_badge.fill.solid()
    l_badge.fill.fore_color.rgb = CORAL_DARK
    l_badge.line.fill.background()
    tf_lb = l_badge.text_frame
    p = tf_lb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "LEVERS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 5 Levers Steps
    levers = [
        ("1", "Wishlist revisit", "Return to saved item in ≤7d"),
        ("2", "Decision resolution", "Fit / comparison doubt is resolved"),
        ("3", "Wishlist → Bag", "Commit the shortlisted item"),
        ("4", "Bag → Purchase", "Complete the transaction"),
        ("5", "30-day window", "Do it before intent decays")
    ]
    w = Inches(2.186)
    gap = Inches(0.2)
    for i, (num, label, desc) in enumerate(levers):
        l_x = Inches(0.8) + i * (w + gap)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x, Inches(3.3), w, Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i != 1 and i != 2 else LIGHT_RED_BG
        card.line.color.rgb = CARD_BORDER if i != 1 and i != 2 else CORAL_DARK

        # Badge inside
        circ = slide2.shapes.add_shape(MSO_SHAPE.OVAL, l_x + Inches(0.15), Inches(3.45), Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE if i != 1 and i != 2 else CORAL_DARK
        circ.line.fill.background()
        tf_c = circ.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE

        tx = slide2.shapes.add_textbox(l_x + Inches(0.55), Inches(3.35), w - Inches(0.6), Inches(1.3))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

    # Bottom Panels
    # Panel 1: Where this project intervenes
    p1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(5.7), Inches(1.8))
    p1.fill.solid()
    p1.fill.fore_color.rgb = LIGHT_BG
    p1.line.color.rgb = CARD_BORDER
    tf1 = p1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.25)
    tf1.margin_top = Inches(0.2)
    p = tf1.paragraphs[0]
    p.text = "Where this project intervenes"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf1.add_paragraph()
    p2.text = "High-intent users already like the item. FitCheck targets the two adjacent, controllable steps: resolve the decision blocker → increase Wishlist-to-Bag progression."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    # Panel 2: What we do NOT optimize
    p2_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(4.9), Inches(5.7), Inches(1.8))
    p2_box.fill.solid()
    p2_box.fill.fore_color.rgb = RGBColor(255, 251, 235) # Light warm yellow
    p2_box.line.color.rgb = GOLD
    tf2 = p2_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.25)
    tf2.margin_top = Inches(0.2)
    p = tf2.paragraphs[0]
    p.text = "What we do NOT optimize"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf2.add_paragraph()
    p2.text = "Generic wishlist adds or time-in-app. More saving is not success unless saved intent converts within 30 days."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: AI Discovery Engine
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "AI-Powered Discovery Engine", "The discovery engine turns 1,500 messy conversations into ranked, traceable opportunities.", "It classifies intent and friction separately, preserves provenance, and exposes the evidence behind every opportunity.")

    pipeline_steps = [
        ("1", "INGEST", "1,500 public conversations\nReddit · Play Store · blogs"),
        ("2", "STRUCTURE", "14-field multi-label schema\nIntent · friction · workaround"),
        ("3", "SEGMENT", "Intent × friction states\nSeparates storage from friction"),
        ("4", "PRIORITIZE", "Opportunity score /25\nP + I + S + A + C"),
        ("5", "AUDIT", "Human review + explorer\n300 checked; links retained")
    ]

    w = Inches(2.186)
    gap = Inches(0.2)
    for i, (num, title, desc) in enumerate(pipeline_steps):
        l_x = Inches(0.8) + i * (w + gap)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x, Inches(2.0), w, Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = CARD_BORDER

        # Circle badge
        circ = slide3.shapes.add_shape(MSO_SHAPE.OVAL, l_x + Inches(0.15), Inches(2.15), Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE
        circ.line.fill.background()
        tf_c = circ.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE

        tx = slide3.shapes.add_textbox(l_x + Inches(0.55), Inches(2.1), w - Inches(0.6), Inches(1.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = CORAL_DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

    # Bottom Left Panel: Evidence Object
    ev = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(7.0), Inches(2.7))
    ev.fill.solid()
    ev.fill.fore_color.rgb = LIGHT_BG
    ev.line.color.rgb = CARD_BORDER
    tf_ev = ev.text_frame
    tf_ev.word_wrap = True
    tf_ev.margin_left = Inches(0.25)
    tf_ev.margin_top = Inches(0.2)
    p = tf_ev.paragraphs[0]
    p.text = "A record becomes a decision-ready evidence object — not just a sentiment label."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Tags inside evidence object
    tags = [
        ("Intent", "HIGH", LIGHT_RED_BG, CORAL_DARK),
        ("Friction", "FIT", RGBColor(254, 243, 199), GOLD),
        ("Workaround", "PAST ORDERS", LIGHT_GREEN_BG, RGBColor(16, 185, 129)),
        ("Source", "REDDIT", LIGHT_BG, NAVY),
        ("Confidence", "0.91", LIGHT_BG, NAVY)
    ]
    for i, (lbl, val, bg_c, txt_c) in enumerate(tags):
        col = i % 3
        row = i // 3
        l = Inches(1.0) + col * Inches(2.1)
        t = Inches(4.6) + row * Inches(0.6)
        tg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(1.9), Inches(0.5))
        tg.fill.solid()
        tg.fill.fore_color.rgb = bg_c
        tg.line.color.rgb = CARD_BORDER
        tf = tg.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"{lbl}: {val}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = txt_c

    p_note = tf_ev.add_paragraph()
    p_note.text = "\nWhy it matters: this lets the PM compare opportunity areas by purchase proximity, severity and addressability — while still drilling back to source evidence."
    p_note.font.size = Pt(14)
    p_note.font.color.rgb = TEXT_MUTED

    # Bottom Right Panel: Quality Control
    qc = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(4.0), Inches(4.433), Inches(2.7))
    qc.fill.solid()
    qc.fill.fore_color.rgb = ACCENT_BLUE
    qc.line.fill.background()
    tf_qc = qc.text_frame
    tf_qc.word_wrap = True
    tf_qc.margin_left = Inches(0.3)
    tf_qc.margin_top = Inches(0.2)
    p = tf_qc.paragraphs[0]
    p.text = "QUALITY CONTROL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p2 = tf_qc.add_paragraph()
    p2.text = "300  records manually reviewed"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    p3 = tf_qc.add_paragraph()
    p3.text = "270 / 300  label agreement = 90%"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    btn_q = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(5.7), Inches(4.033), Inches(0.6))
    btn_q.fill.solid()
    btn_q.fill.fore_color.rgb = CORAL_DARK
    btn_q.line.fill.background()
    tf_bq = btn_q.text_frame
    p = tf_bq.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Test the engine live 🔗"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 4: Discovery Findings
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "Discovery Findings", "Fit uncertainty is the highest-value non-monetary wedge — not simply the most frequent complaint.", "The engine separates purchase intent from friction, then scores opportunities on prevalence, intent proximity, severity, addressability and confidence.")

    # 2x2 Grid (Left Side)
    grid_title = slide4.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(0.35))
    tf_gt = grid_title.text_frame
    p = tf_gt.paragraphs[0]
    p.text = "Intent × friction distribution (n = 1,500)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    cards_2x2 = [
        ("51.3%", "HIGH INTENT / HIGH FRICTION\nPrimary growth target · n=770", CORAL_DARK, WHITE, Inches(0.8), Inches(2.3)),
        ("22.0%", "LOW INTENT / HIGH FRICTION\nPassive browsing · n=330", RGBColor(254, 243, 199), NAVY, Inches(3.8), Inches(2.3)),
        ("16.7%", "HIGH INTENT / LOW FRICTION\nWaiting to execute · n=250", LIGHT_GREEN_BG, NAVY, Inches(0.8), Inches(4.4)),
        ("10.0%", "LOW INTENT / LOW FRICTION\nMood-board archive · n=150", LIGHT_BG, NAVY, Inches(3.8), Inches(4.4))
    ]

    for pct, desc, bg_c, txt_c, l, t in cards_2x2:
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(2.8), Inches(1.9))
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = CARD_BORDER if bg_c != CORAL_DARK else CORAL_DARK
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = pct
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = txt_c
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = txt_c

    # Right Side: Opportunity Ranking Bars
    rank_title = slide4.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.35))
    tf_rt = rank_title.text_frame
    p = tf_rt.paragraphs[0]
    p.text = "Opportunity ranking (/25)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    rankings = [
        ("1 Size & Fit Uncertainty", 23, CORAL_DARK),
        ("2 Comparison Paralysis", 22, ACCENT_BLUE),
        ("3 Product Reality / Quality", 19, RGBColor(13, 148, 136)),
        ("4 Price Volatility & Context", 17, GOLD),
        ("5 Styling & Wardrobe Match", 16, RGBColor(100, 116, 139))
    ]

    for i, (label, score, bar_c) in enumerate(rankings):
        t_y = Inches(2.4) + i * Inches(0.7)
        # Label
        tx = slide4.shapes.add_textbox(Inches(7.0), t_y, Inches(3.2), Inches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Progress bar background
        bar_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), t_y + Inches(0.08), Inches(1.8), Inches(0.25))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_BG
        bar_bg.line.fill.background()

        # Progress bar fill
        fill_w = Inches(1.8) * (score / 25.0)
        bar_fill = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), t_y + Inches(0.08), fill_w, Inches(0.25))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = bar_c
        bar_fill.line.fill.background()

        # Score text
        tx_s = slide4.shapes.add_textbox(Inches(12.1), t_y, Inches(0.6), Inches(0.4))
        tf_s = tx_s.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = str(score)
        p_s.font.size = Pt(14)
        p_s.font.bold = True
        p_s.font.color.rgb = NAVY

    # Callout Box at bottom right
    callout = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.8), Inches(5.533), Inches(0.9))
    callout.fill.solid()
    callout.fill.fore_color.rgb = LIGHT_RED_BG
    callout.line.color.rgb = CORAL_DARK
    tf_co = callout.text_frame
    tf_co.word_wrap = True
    tf_co.margin_left = Inches(0.2)
    tf_co.margin_top = Inches(0.15)
    p = tf_co.paragraphs[0]
    p.text = "Frequency ≠ priority. Fit wins because the affected users are closer to purchase, the friction is severe, and the problem is addressable without monetary incentives."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CORAL_DARK

    # =========================================================================
    # SLIDE 5: Primary Research
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "Primary Research", "Interviews narrowed the hypothesis: fit appeared in 5/6 journeys and was the primary blocker for 3/6.", "Six wishlist walkthroughs reconstructed the path from “saved” to “researched” to “delayed / bought” for near-term purchase intent.")

    # 3 Stat Cards Top
    stats5 = [
        ("5 / 6", "mentioned fit or drape somewhere in the decision journey"),
        ("3 / 6", "identified fit / size as the primary blocker"),
        ("6 / 6", "used an outside-app workaround before deciding")
    ]
    for i, (num, desc) in enumerate(stats5):
        l = Inches(0.8) + i * Inches(3.98)
        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, Inches(1.9), Inches(3.78), Inches(1.1))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_BG if i != 1 else LIGHT_RED_BG
        c.line.color.rgb = CARD_BORDER if i != 1 else CORAL_DARK

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = NAVY if i != 1 else CORAL_DARK

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

    # Left: Interview table
    t_title = slide5.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(5.8), Inches(0.35))
    tf_tt = t_title.text_frame
    p = tf_tt.paragraphs[0]
    p.text = "What each participant was waiting on"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    rows, cols = 7, 3
    table5 = slide5.shapes.add_table(rows, cols, Inches(0.8), Inches(3.5), Inches(5.8), Inches(3.0)).table
    table5.columns[0].width = Inches(0.8)
    table5.columns[1].width = Inches(2.5)
    table5.columns[2].width = Inches(2.5)

    headers = ["P#", "Wishlisted item", "Primary blocker"]
    for j, h in enumerate(headers):
        cell = table5.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    p_data = [
        ("P1", "Roadster jeans", "Fit / size"),
        ("P2", "Anouk dress", "Fit & drape"),
        ("P3", "HRX shoes", "Comparison"),
        ("P4", "Libas kurti set", "Post-wash fit"),
        ("P5", "Allen Solly blazer", "Fabric quality"),
        ("P6", "Highlander cargos", "Styling")
    ]

    for i, row in enumerate(p_data):
        for j, val in enumerate(row):
            cell = table5.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            if j == 2 and ("Fit" in val or "fit" in val):
                p.font.bold = True
                p.font.color.rgb = CORAL_DARK

    # Right: Observed workaround loop
    w_title = slide5.shapes.add_textbox(Inches(7.0), Inches(3.1), Inches(5.5), Inches(0.35))
    tf_wt = w_title.text_frame
    p = tf_wt.paragraphs[0]
    p.text = "Observed workaround loop"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    loop_steps = [
        ("1", "SAVE", "Specific item / size"),
        ("2", "LEAVE APP", "Reddit · YouTube · Instagram"),
        ("3", "TRIANGULATE", "Past orders · reviews · social proof"),
        ("4", "DELAY", "Still not confident enough to commit")
    ]

    for i, (num, lbl, desc) in enumerate(loop_steps):
        t_y = Inches(3.5) + i * Inches(0.65)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), t_y, Inches(5.533), Inches(0.58))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG if i != 1 else LIGHT_RED_BG
        box.line.color.rgb = CARD_BORDER if i != 1 else CORAL_DARK

        # Number circle
        circ = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), t_y + Inches(0.1), Inches(0.38), Inches(0.38))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE if i != 1 else CORAL_DARK
        circ.line.fill.background()
        tf_c = circ.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE

        tx = slide5.shapes.add_textbox(Inches(7.6), t_y, Inches(4.8), Inches(0.58))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = f"{lbl}:  {desc}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

    # Verbatim Quote Banner at bottom right
    q_b = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(6.15), Inches(5.533), Inches(0.55))
    q_b.fill.solid()
    q_b.fill.fore_color.rgb = LIGHT_RED_BG
    q_b.line.color.rgb = CORAL_DARK
    tf_qb = q_b.text_frame
    tf_qb.margin_left = Inches(0.2)
    p = tf_qb.paragraphs[0]
    p.text = "“Roadster 30 is too tight; 32 slips off. I compared three previous orders.” — P1"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = CORAL_DARK

    # =========================================================================
    # SLIDE 6: Problem Definition
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "Problem Definition", "The root problem is not lack of desire — it is lack of trusted, personalized fit confidence at the commitment moment.", "The thinking narrows consistently from the business metric to a specific user, moment, root cause and product outcome.")

    # Top Flow Steps (Funnel Evolution)
    flow_steps = [
        ("BUSINESS METRIC", "30D wishlist buyers"),
        ("PRODUCT OUTCOME", "Wishlist → Bag"),
        ("AI DISCOVERY", "Fit #1 · 23/25"),
        ("INTERVIEWS", "5/6 surfaced fit"),
        ("ROOT PROBLEM", "Cross-brand fit trust")
    ]
    w = Inches(2.186)
    gap = Inches(0.2)
    for i, (lbl, val) in enumerate(flow_steps):
        l_x = Inches(0.8) + i * (w + gap)
        c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x, Inches(1.85), w, Inches(0.8))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_BG if i != 4 else LIGHT_RED_BG
        c.line.color.rgb = CARD_BORDER if i != 4 else CORAL_DARK
        tf = c.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED if i != 4 else CORAL_DARK
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = NAVY if i != 4 else CORAL_DARK

    # Problem Statement Central Banner
    ps_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(11.733), Inches(1.3))
    ps_box.fill.solid()
    ps_box.fill.fore_color.rgb = ACCENT_BLUE
    ps_box.line.fill.background()
    tf_ps = ps_box.text_frame
    tf_ps.word_wrap = True
    tf_ps.margin_left = Inches(0.3)
    tf_ps.margin_top = Inches(0.2)
    p = tf_ps.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "High-intent apparel shoppers with a near-term purchase goal delay moving a wishlisted item to Bag because they cannot confidently predict how their selected size will fit across brands — so they leave Myntra to triangulate fit from reviews, past orders and try-on content before committing."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 6 Aspect Cards (2x3 Grid)
    aspects = [
        ("WHO", "High-intent apparel shoppers with a near-term purchase goal"),
        ("WHEN", "Revisiting a specific saved SKU to decide “buy now or wait?”"),
        ("ROOT CAUSE", "Size labels are not trusted as personal, cross-brand fit evidence"),
        ("WORKAROUND", "Reddit / YouTube / Instagram + comparing previous kept orders"),
        ("USER VALUE", "Decide fit without leaving Myntra or ordering “just to try”"),
        ("BUSINESS VALUE", "Increase high-intent Wishlist → Bag; fit accuracy may also reduce returns")
    ]

    for i, (title, desc) in enumerate(aspects):
        col = i % 3
        row = i // 3
        l = Inches(0.8) + col * Inches(3.98)
        t = Inches(4.3) + row * Inches(1.3)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(3.78), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = CARD_BORDER
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = CORAL_DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MAIN

    # =========================================================================
    # SLIDE 7: Solution Rationale
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "Solution Rationale", "Smart Wishlist + FitCheck resolves the declared blocker inside the wishlist, then hands control back to the user.", "The core MVP is a decision-confidence loop — not a discount, coupon, or forced checkout nudge.")

    # 5 Step Mechanism Flow
    steps7 = [
        ("1", "FIT ANCHOR", "Confirm a kept size\nPast fit = reference", "Roadster · M"),
        ("2", "WAITING ON FIT", "Declare the blocker\nCaptures why waiting", "Waiting: Fit"),
        ("3", "FITCHECK", "Explain + recommend\nReturn size signal", "High confidence"),
        ("4", "READY TO BUY", "Resolve condition\nFlips state when ready", "READY TO BUY"),
        ("5", "MOVE TO BAG", "User commits\nShopper takes action", "MOVE TO BAG")
    ]

    w = Inches(2.186)
    gap = Inches(0.2)
    for i, (num, label, desc, pill) in enumerate(steps7):
        l_x = Inches(0.8) + i * (w + gap)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x, Inches(2.0), w, Inches(2.7))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i != 3 else LIGHT_RED_BG
        card.line.color.rgb = CARD_BORDER if i != 3 else CORAL_DARK

        # Circle badge
        circ = slide7.shapes.add_shape(MSO_SHAPE.OVAL, l_x + Inches(0.15), Inches(2.15), Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT_BLUE if i != 3 else CORAL_DARK
        circ.line.fill.background()
        tf_c = circ.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        p_c.text = num
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = WHITE

        tx = slide7.shapes.add_textbox(l_x + Inches(0.55), Inches(2.1), w - Inches(0.6), Inches(1.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED

        # Bottom Pill inside card
        pill_b = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x + Inches(0.15), Inches(4.1), w - Inches(0.3), Inches(0.45))
        pill_b.fill.solid()
        pill_b.fill.fore_color.rgb = WHITE if i != 3 else CORAL_DARK
        pill_b.line.color.rgb = CARD_BORDER if i != 3 else CORAL_DARK
        tf_p = pill_b.text_frame
        p_p = tf_p.paragraphs[0]
        p_p.alignment = PP_ALIGN.CENTER
        p_p.text = pill
        p_p.font.size = Pt(14)
        p_p.font.bold = True
        p_p.font.color.rgb = CORAL_DARK if i != 3 else WHITE

    # Bottom Panels
    # Panel 1: Why this fits brief
    p1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(7.0), Inches(1.4))
    p1.fill.solid()
    p1.fill.fore_color.rgb = LIGHT_GREEN_BG
    p1.line.color.rgb = RGBColor(16, 185, 129)
    tf1 = p1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.25)
    tf1.margin_top = Inches(0.15)
    p = tf1.paragraphs[0]
    p.text = "Why this fits the brief"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)
    p2 = tf1.add_paragraph()
    p2.text = "No subsidy. No coupon. The value is informational: translate fit evidence into a personal recommendation, then let the shopper decide."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MAIN

    # Panel 2: Scope discipline
    p2_b = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(4.9), Inches(4.433), Inches(1.4))
    p2_b.fill.solid()
    p2_b.fill.fore_color.rgb = RGBColor(255, 251, 235)
    p2_b.line.color.rgb = GOLD
    tf2 = p2_b.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.25)
    tf2.margin_top = Inches(0.15)
    p = tf2.paragraphs[0]
    p.text = "Scope discipline"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p2 = tf2.add_paragraph()
    p2.text = "Price / Both are expansion modes in the prototype. The causal experiment isolates Fit-only."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MAIN

    # Button
    btn_m = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(6.45), Inches(3.533), Inches(0.55))
    btn_m.fill.solid()
    btn_m.fill.fore_color.rgb = CORAL_DARK
    btn_m.line.fill.background()
    tf_bm = btn_m.text_frame
    p = tf_bm.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Launch live MVP 🚀"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # =========================================================================
    # SLIDE 8: Deployed MVP Scope
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide8, "Deployed MVP", "The deployed MVP is intentionally thin: it proves the decision-resolution loop before investing in real fit ML.", "A top-scoring MVP is explicit about what is functional, what is simulated, and exactly what learning it unlocks.")

    # Left Box: Functional in live prototype
    f_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.7), Inches(3.8))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = LIGHT_BG
    f_box.line.color.rgb = CARD_BORDER
    tf_f = f_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = Inches(0.3)
    tf_f.margin_top = Inches(0.2)
    p = tf_f.paragraphs[0]
    p.text = "Functional in the live prototype"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY

    funcs = [
        "✓  Set and save a Fit Anchor profile",
        "✓  Declare the blocker: Fit / Price / Both",
        "✓  Run a FitCheck decision event",
        "✓  Transition the item to “Ready to Buy”",
        "✓  Move the resolved item to Bag"
    ]
    for f in funcs:
        p_f = tf_f.add_paragraph()
        p_f.text = f
        p_f.font.size = Pt(14)
        p_f.font.color.rgb = TEXT_MAIN

    p_demo = tf_f.add_paragraph()
    p_demo.text = "\nDemo extension: simulated price events for Price / Both states."
    p_demo.font.size = Pt(14)
    p_demo.font.italic = True
    p_demo.font.color.rgb = TEXT_MUTED

    btn_t = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(4.9), Inches(2.5), Inches(0.55))
    btn_t.fill.solid()
    btn_t.fill.fore_color.rgb = CORAL_DARK
    btn_t.line.fill.background()
    tf_bt = btn_t.text_frame
    p = tf_bt.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Test live MVP 🚀"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Right Box: Not production-ready yet
    nf_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.9), Inches(5.7), Inches(3.8))
    nf_box.fill.solid()
    nf_box.fill.fore_color.rgb = ACCENT_BLUE
    nf_box.line.fill.background()
    tf_nf = nf_box.text_frame
    tf_nf.word_wrap = True
    tf_nf.margin_left = Inches(0.3)
    tf_nf.margin_top = Inches(0.2)
    p = tf_nf.paragraphs[0]
    p.text = "Not production-ready yet — by design"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    nfuncs = [
        "•  Real order / return-history ingestion",
        "•  Production fit model + category calibration",
        "•  Live inventory / price event integrations",
        "•  Notification preferences + frequency caps",
        "•  Learned confidence thresholds + abstention logic"
    ]
    for nf in nfuncs:
        p_nf = tf_nf.add_paragraph()
        p_nf.text = nf
        p_nf.font.size = Pt(14)
        p_nf.font.color.rgb = RGBColor(226, 232, 240)

    p_trans = tf_nf.add_paragraph()
    p_trans.text = "\nTransparency prevents “prototype theatre”: this MVP validates the state-transition mechanism, not model accuracy."
    p_trans.font.size = Pt(14)
    p_trans.font.bold = True
    p_trans.font.color.rgb = GOLD

    # Bottom 60-Second Evaluator Path
    eval_path = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.85), Inches(11.733), Inches(0.9))
    eval_path.fill.solid()
    eval_path.fill.fore_color.rgb = LIGHT_BG
    eval_path.line.color.rgb = CARD_BORDER
    tf_ep = eval_path.text_frame
    tf_ep.word_wrap = True
    tf_ep.margin_left = Inches(0.2)
    tf_ep.margin_top = Inches(0.2)
    p = tf_ep.paragraphs[0]
    p.text = "60-second evaluator path:   [1] Set anchor  ──>  [2] Choose Fit  ──>  [3] Run FitCheck  ──>  [4] See Ready  ──>  [5] Move to Bag"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # =========================================================================
    # SLIDE 9: Define Success
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide9, "Define Success", "Success is a causal lift in 30-day wishlist buyers, with faster fit-to-bag progression and no return penalty.", "The metric stack separates business impact, mechanism validation, leading indicators and guardrails — each with a clear decision use.")

    # 3 Top Metric Cards
    metrics_top = [
        ("BUSINESS METRIC", "30-Day Wishlist Buyer Conversion", "% unique wishlist users who purchase ≥1 saved item within 30 days of the add.", ACCENT_BLUE, WHITE),
        ("PRIMARY MECHANISM", "Fit-Resolved → Bag within 72h", "Of fit-tagged items receiving a decisive FitCheck, % moved to Bag within 72 hours.", LIGHT_RED_BG, CORAL_DARK),
        ("KEY GUARDRAIL", "Fit-related return rate", "Must not increase. Also track notification opt-out so confidence nudges do not erode trust.", LIGHT_GREEN_BG, RGBColor(16, 185, 129))
    ]

    w = Inches(3.75)
    gap = Inches(0.24)
    for i, (cat, title, desc, bg_c, txt_c) in enumerate(metrics_top):
        l_x = Inches(0.8) + i * (w + gap)
        c = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_x, Inches(1.9), w, Inches(1.7))
        c.fill.solid()
        c.fill.fore_color.rgb = bg_c
        c.line.color.rgb = CARD_BORDER if bg_c != LIGHT_RED_BG else CORAL_DARK
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = txt_c if bg_c == ACCENT_BLUE else CORAL_DARK
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = txt_c if bg_c == ACCENT_BLUE else NAVY
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = txt_c if bg_c == ACCENT_BLUE else TEXT_MUTED

    # Leading Indicators Table
    t_title = slide9.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.733), Inches(0.35))
    tf_tt = t_title.text_frame
    p = tf_tt.paragraphs[0]
    p.text = "Leading indicators + guardrails"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY

    rows, cols = 6, 3
    table9 = slide9.shapes.add_table(rows, cols, Inches(0.8), Inches(4.1), Inches(11.733), Inches(2.1)).table
    table9.columns[0].width = Inches(2.733)
    table9.columns[1].width = Inches(4.5)
    table9.columns[2].width = Inches(4.5)

    headers = ["Metric", "Definition", "Why it matters"]
    for j, h in enumerate(headers):
        cell = table9.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    m_data = [
        ("Anchor adoption", "% fit-tagged users with a usable fit anchor", "Can users reach the value moment?"),
        ("Fit resolution rate", "% FitChecks returning a decisive, explainable recommendation", "Does the system actually resolve uncertainty?"),
        ("Median wishlist → Bag time", "Hours from wishlist add / revisit to Bag for fit cohort", "Is decision latency shrinking?"),
        ("Wishlist add rate", "Eligible users adding ≥1 item to wishlist", "Ensure extra friction does not suppress saving"),
        ("Recommendation disagreement", "% users who override / reject the suggested size", "Early trust and calibration signal")
    ]

    for i, row in enumerate(m_data):
        for j, val in enumerate(row):
            cell = table9.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            if j == 0:
                p.font.bold = True

    # Bottom Experiment Banner
    exp = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.35), Inches(11.733), Inches(0.55))
    exp.fill.solid()
    exp.fill.fore_color.rgb = LIGHT_RED_BG
    exp.line.color.rgb = CORAL_DARK
    tf_e = exp.text_frame
    tf_e.margin_left = Inches(0.2)
    p = tf_e.paragraphs[0]
    p.text = "Experiment: user-level A/B · control=current wishlist · treatment=FitCheck · 30-day observation · segment by category / brand · launch only with reliable business uplift."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CORAL_DARK

    # =========================================================================
    # SLIDE 10: Risks & Mitigation
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide10, "Risks & Mitigation", "The biggest failure mode is false confidence; rollout should privilege abstention, trust and learning over reach.", "The solution fails if it pushes users faster into wrong-size orders, creates setup friction, or generalizes a narrow research sample too aggressively.")

    # Table of Risks
    rows, cols = 6, 4
    table10 = slide10.shapes.add_table(rows, cols, Inches(0.8), Inches(1.9), Inches(11.733), Inches(4.3)).table
    table10.columns[0].width = Inches(2.5)
    table10.columns[1].width = Inches(1.0)
    table10.columns[2].width = Inches(3.8)
    table10.columns[3].width = Inches(4.433)

    headers = ["Risk", "Impact", "Why it could happen", "Mitigation / rollout gate"]
    for j, h in enumerate(headers):
        cell = table10.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    risk_data = [
        ("False fit confidence → returns", "HIGH", "Prototype reasoning can be overconfident without garment measurement truth", "Confidence bands + “not enough evidence” state · pilot stable categories · calibrate on kept/returned orders"),
        ("Fit Anchor setup friction", "MED", "Manual multi-field setup delays value", "Auto-suggest from kept orders; ask for one-tap confirmation contextually, not upfront"),
        ("Cold start / sparse brand data", "HIGH", "New brands, categories or unusual body / size combinations lack evidence", "Blend size charts + material + brand signals; abstain below threshold instead of guessing"),
        ("Research sample does not generalize", "MED", "6 interviews are directional; public-conversation mix may not equal Myntra population", "User-level A/B before rollout; report effects by category, brand, gender and new vs repeat users"),
        ("Notification fatigue / trust erosion", "MED", "“Ready to Buy” nudges can feel pushy or repetitive", "Trigger only on meaningful condition change · frequency cap · digest · easy opt-out")
    ]

    for i, row in enumerate(risk_data):
        for j, val in enumerate(row):
            cell = table10.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MAIN
            if j == 0:
                p.font.bold = True
            elif j == 1:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.color.rgb = CORAL_DARK if val == "HIGH" else GOLD

    # Bottom Rollout Principle Box
    r_princ = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.6))
    r_princ.fill.solid()
    r_princ.fill.fore_color.rgb = ACCENT_BLUE
    r_princ.line.fill.background()
    tf_rp = r_princ.text_frame
    tf_rp.margin_left = Inches(0.2)
    tf_rp.margin_top = Inches(0.12)
    p = tf_rp.paragraphs[0]
    p.text = "Rollout principle:  Earn trust before reach: offline calibration → Fit-only pilot → category expansion.  “A wishlist should manage the decision, not just remember the product.”"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE

    output_path = "Myntra_Wishlist_To_Purchase_Growth.pptx"
    prs.save(output_path)
    print(f"Successfully generated updated PDF-matched presentation at {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
