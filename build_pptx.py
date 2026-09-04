import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # Color Palette (Color-blind friendly & high contrast)
    BG_OFFWHITE = RGBColor(248, 250, 252) # #f8fafc
    WHITE = RGBColor(255, 255, 255)
    NAVY_DARK = RGBColor(15, 23, 42) # #0f172a
    NAVY_CARD = RGBColor(27, 38, 79) # #1b264f
    ROSE = RGBColor(225, 29, 72) # #e11d48
    ROSE_BG = RGBColor(255, 241, 242) # #fff1f2
    AMBER = RGBColor(217, 119, 6) # #d97706
    AMBER_BG = RGBColor(254, 243, 199) # #fef3c7
    AMBER_CARD_BG = RGBColor(255, 251, 235) # #fffbeb
    GREEN = RGBColor(16, 185, 129) # #10b981
    GREEN_BG = RGBColor(236, 253, 245) # #ecfdf5
    SLATE_TEXT = RGBColor(51, 65, 85) # #334155
    MUTED_TEXT = RGBColor(71, 85, 105) # #475569
    LIGHT_MUTED = RGBColor(100, 116, 139) # #64748b
    BORDER_COLOR = RGBColor(203, 213, 225) # #cbd5e1

    FONT_FAMILY = "Inter"

    def add_bg(slide, color=BG_OFFWHITE):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, kicker_text, title_text, subtitle_text=None, top_start=0.5):
        tx_box = slide.shapes.add_textbox(Inches(0.7), Inches(top_start), Inches(14.6), Inches(1.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Kicker
        p0 = tf.paragraphs[0]
        p0.text = kicker_text.upper()
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(14) # Minimum 14pt rule
        p0.font.bold = True
        p0.font.color.rgb = ROSE
        p0.space_after = Pt(4)

        # Action Title
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(4)

        current_height = 0.8
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(14) # Minimum 14pt rule
            p2.font.color.rgb = MUTED_TEXT
            current_height = 1.2
        return top_start + current_height

    def add_footer(slide, slide_num):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(8.3), Inches(14.6), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_COLOR
        line.line.fill.background()

        tx_left = slide.shapes.add_textbox(Inches(0.7), Inches(8.35), Inches(10), Inches(0.4))
        tf_left = tx_left.text_frame
        tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0
        p_l = tf_left.paragraphs[0]
        p_l.text = "Myntra · Wishlist-to-Purchase Growth Case Study"
        p_l.font.name = FONT_FAMILY
        p_l.font.size = Pt(14)
        p_l.font.color.rgb = MUTED_TEXT

        tx_right = slide.shapes.add_textbox(Inches(12.3), Inches(8.35), Inches(3.0), Inches(0.4))
        tf_right = tx_right.text_frame
        tf_right.margin_left = tf_right.margin_top = tf_right.margin_right = tf_right.margin_bottom = 0
        p_r = tf_right.paragraphs[0]
        p_r.text = f"Slide {slide_num}"
        p_r.alignment = PP_ALIGN.RIGHT
        p_r.font.name = FONT_FAMILY
        p_r.font.size = Pt(14)
        p_r.font.bold = True
        p_r.font.color.rgb = MUTED_TEXT

    def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=None, left_accent_width=0.08):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()

        if left_accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(left_accent_width), Inches(height))
            accent.fill.solid()
            accent.fill.fore_color.rgb = left_accent_color
            accent.line.fill.background()
        return shape

    def add_button(slide, left, top, width, height, text, url, bg_color=NAVY_DARK):
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        btn.fill.solid()
        btn.fill.fore_color.rgb = bg_color
        btn.line.fill.background()

        tf = btn.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14) # Minimum 14pt rule
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        if url:
            btn.click_action.hyperlink.address = url
        return btn

    def add_table(slide, left, top, width, height, headers, data, col_widths=None):
        rows = len(data) + 1
        cols = len(headers)
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table

        if col_widths:
            for idx, w in enumerate(col_widths):
                table.columns[idx].width = Inches(w)

        # Header Row
        for c_idx, h_text in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_CARD
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = Inches(0.08)
            p = tf.paragraphs[0]
            p.text = h_text
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE

        # Data Rows
        for r_idx, row_data in enumerate(data):
            row_bg = WHITE if r_idx % 2 == 1 else BG_OFFWHITE
            for c_idx, val in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                tf = cell.text_frame
                tf.margin_left = tf.margin_right = Inches(0.12)
                tf.margin_top = tf.margin_bottom = Inches(0.08)
                p = tf.paragraphs[0]
                
                if isinstance(val, tuple):
                    p.text = str(val[0])
                    p.font.bold = val[1]
                    if len(val) > 2 and val[2]:
                        p.font.color.rgb = val[2]
                    else:
                        p.font.color.rgb = NAVY_DARK
                else:
                    p.text = str(val)
                    p.font.bold = False
                    p.font.color.rgb = NAVY_DARK

                p.font.name = FONT_FAMILY
                p.font.size = Pt(14)

        return table_shape

    # ==========================================
    # SLIDE 1: Cover / Executive Summary
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    tx_box1 = slide1.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.5), Inches(7.5))
    tf1 = tx_box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p = tf1.paragraphs[0]
    p.text = "GROWTH PM CASE STUDY · MYNTRA"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ROSE
    p.space_after = Pt(10)

    p = tf1.add_paragraph()
    p.text = "Convert more wishlisters by resolving the last uncertainty — not by discounting."
    p.font.name = FONT_FAMILY
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_DARK
    p.space_after = Pt(14)

    p = tf1.add_paragraph()
    p.text = "Smart Wishlist + FitCheck turns a passive save into a decision-ready state for high-intent apparel shoppers — and measures whether that confidence actually moves them to Bag and purchase within 30 days."
    p.font.name = FONT_FAMILY
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED_TEXT

    # 3 Stat Cards
    add_card(slide1, 0.7, 4.3, 2.7, 1.5, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx = slide1.shapes.add_textbox(Inches(0.85), Inches(4.4), Inches(2.4), Inches(1.3))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "51.3%"; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = ROSE
    p = tf.add_paragraph(); p.text = "of analyzed conversations fell in high-intent / high-friction"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    add_card(slide1, 3.6, 4.3, 2.7, 1.5, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=AMBER)
    tx = slide1.shapes.add_textbox(Inches(3.75), Inches(4.4), Inches(2.4), Inches(1.3))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "23/25"; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = AMBER
    p = tf.add_paragraph(); p.text = "Fit uncertainty ranked #1 on opportunity rubric"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    add_card(slide1, 6.5, 4.3, 2.7, 1.5, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=GREEN)
    tx = slide1.shapes.add_textbox(Inches(6.65), Inches(4.4), Inches(2.4), Inches(1.3))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "5 / 6"; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = GREEN
    p = tf.add_paragraph(); p.text = "interviews surfaced fit / drape uncertainty"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    # Action Buttons
    add_button(slide1, 0.7, 6.1, 3.2, 0.65, "Open AI Discovery Engine 🔗", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/?v=10", bg_color=NAVY_DARK)
    add_button(slide1, 4.1, 6.1, 2.5, 0.65, "Open Live MVP 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    tx_note = slide1.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(8.5), Inches(0.5))
    p = tx_note.text_frame.paragraphs[0]
    p.text = "No monetary incentives · 10-slide submission · live artefacts linked"
    p.font.size = Pt(14); p.font.color.rgb = LIGHT_MUTED

    # Right Panel: THE MECHANISM
    add_card(slide1, 9.7, 0.6, 5.6, 7.3, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=ROSE, left_accent_width=0.1)
    tx_mech = slide1.shapes.add_textbox(Inches(10.0), Inches(0.9), Inches(5.0), Inches(6.7))
    tf_mech = tx_mech.text_frame; tf_mech.word_wrap = True

    p = tf_mech.paragraphs[0]
    p.text = "THE MECHANISM"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(6)

    p = tf_mech.add_paragraph()
    p.text = "From “Saved” to “Ready to Buy”"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(16)

    steps = [
        ("1. SAVED: ", "User has explicit interest, not commitment"),
        ("2. FIT UNCERTAIN: ", "Cross-brand size / drape is unresolved"),
        ("3. FIT RESOLVED: ", "FitCheck returns transparent confidence signal"),
        ("4. READY TO BUY: ", "User decides; item becomes action-ready"),
        ("5. BAG → PURCHASE: ", "Commitment advances toward 30-day goal")
    ]
    for s_title, s_desc in steps:
        p = tf_mech.add_paragraph()
        run1 = p.add_run(); run1.text = s_title; run1.font.bold = True; run1.font.color.rgb = ROSE; run1.font.size = Pt(14)
        run2 = p.add_run(); run2.text = s_desc; run2.font.color.rgb = SLATE_TEXT; run2.font.size = Pt(14)
        p.space_after = Pt(10)

    add_card(slide1, 10.0, 6.2, 5.0, 1.2, bg_color=ROSE_BG, border_color=ROSE)
    tx_h = slide1.shapes.add_textbox(Inches(10.1), Inches(6.3), Inches(4.8), Inches(1.0))
    tf_h = tx_h.text_frame; tf_h.word_wrap = True
    p = tf_h.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thesis: decision confidence is the addressable growth lever."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    add_footer(slide1, 1)

    # ==========================================
    # SLIDE 2: Business Metric Decomposition
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "BUSINESS METRIC DECOMPOSITION", 
               "The business metric moves only when wishlisted users progress from intent to commitment.",
               "The brief is user-level, so the North Star must also be user-level — not an SKU conversion proxy.")

    add_card(slide2, 0.7, 1.9, 14.6, 1.0, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER, left_accent_width=0.08)
    tx_f = slide2.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(14.2), Inches(0.8))
    tf_f = tx_f.text_frame; tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "30-Day Wishlist Buyer Conversion = (Unique users who purchase ≥1 item they wishlisted within 30 days of that add) ÷ (Unique users with ≥1 eligible wishlist add in cohort window)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE

    funnel_items = [
        ("1. Wishlist Revisit", "Return to saved item in ≤7d", False),
        ("2. Decision Resolution", "Fit / comparison doubt is resolved", True),
        ("3. Wishlist → Bag", "Commit the shortlisted item", True),
        ("4. Bag → Purchase", "Complete the transaction", False),
        ("5. 30-Day Window", "Do it before intent decays", False),
    ]
    card_w = 2.76
    for idx, (f_title, f_desc, is_hi) in enumerate(funnel_items):
        c_left = 0.7 + idx * (card_w + 0.2)
        bg_c = ROSE_BG if is_hi else BG_OFFWHITE
        brd_c = ROSE if is_hi else BORDER_COLOR
        add_card(slide2, c_left, 3.1, card_w, 2.0, bg_color=bg_c, border_color=brd_c)
        tx_box = slide2.shapes.add_textbox(Inches(c_left + 0.15), Inches(3.25), Inches(card_w - 0.3), Inches(1.7))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f_title
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE if is_hi else NAVY_DARK
        p.space_after = Pt(6)
        p = tf.add_paragraph()
        p.text = f_desc
        p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide2, 0.7, 5.3, 7.1, 2.7, bg_color=WHITE, border_color=BORDER_COLOR)
    tx = slide2.shapes.add_textbox(Inches(0.9), Inches(5.5), Inches(6.7), Inches(2.3))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Where this project intervenes"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "High-intent users already like the item. FitCheck targets the two adjacent, controllable steps: resolve the decision blocker → increase Wishlist-to-Bag progression."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide2, 8.2, 5.3, 7.1, 2.7, bg_color=AMBER_CARD_BG, border_color=AMBER)
    tx = slide2.shapes.add_textbox(Inches(8.4), Inches(5.5), Inches(6.7), Inches(2.3))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "What we do NOT optimize"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "Generic wishlist adds or time-in-app. More saving is not success unless saved intent converts within 30 days."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_footer(slide2, 2)

    # ==========================================
    # SLIDE 3: AI-Powered Discovery Engine
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "AI-POWERED DISCOVERY ENGINE",
               "The discovery engine turns 1,500 messy conversations into ranked, traceable opportunities.",
               "It classifies intent and friction separately, preserves provenance, and exposes the evidence behind every opportunity.")

    proc_steps = [
        ("1. INGEST", "1,500 public conversations"),
        ("2. STRUCTURE", "14-field multi-label schema"),
        ("3. SEGMENT", "Intent × friction states"),
        ("4. PRIORITIZE", "Opportunity score /25"),
        ("5. AUDIT", "Human review + explorer"),
    ]
    for idx, (p_title, p_desc) in enumerate(proc_steps):
        c_left = 0.7 + idx * (card_w + 0.2)
        add_card(slide3, c_left, 1.9, card_w, 1.6, bg_color=WHITE, border_color=BORDER_COLOR)
        tx_box = slide3.shapes.add_textbox(Inches(c_left + 0.15), Inches(2.05), Inches(card_w - 0.3), Inches(1.3))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(6)
        p = tf.add_paragraph()
        p.text = p_desc
        p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide3, 0.7, 3.7, 8.8, 4.3, bg_color=WHITE, border_color=BORDER_COLOR)
    tx = slide3.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(8.2), Inches(3.9))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A record becomes a decision-ready evidence object — not just a sentiment label."
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(12)

    p = tf.add_paragraph(); p.space_after = Pt(14)
    run = p.add_run(); run.text = " Intent: HIGH "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = ROSE
    run = p.add_run(); run.text = "   Friction: FIT "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = AMBER
    run = p.add_run(); run.text = "   Workaround: PAST ORDERS "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = GREEN

    p = tf.add_paragraph()
    p.text = "Why it matters: this lets the PM compare opportunity areas by purchase proximity, severity and addressability — while still drilling back to source evidence."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide3, 9.8, 3.7, 5.5, 4.3, bg_color=NAVY_CARD, border_color=None)
    tx = slide3.shapes.add_textbox(Inches(10.0), Inches(3.9), Inches(5.1), Inches(3.7))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "QUALITY CONTROL"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(8)

    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    p.text = "300 records manually reviewed"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.space_after = Pt(8)

    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    p.text = "270 / 300 label agreement = 90%"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BG_OFFWHITE; p.space_after = Pt(18)

    add_button(slide3, 10.8, 6.6, 3.5, 0.65, "Test the engine live 🔗", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/?v=10", bg_color=ROSE)

    add_footer(slide3, 3)

    # ==========================================
    # SLIDE 4: Discovery Findings
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "DISCOVERY FINDINGS",
               "Fit uncertainty is the highest-value non-monetary wedge — not simply the most frequent complaint.",
               "The engine separates purchase intent from friction, then scores opportunities on prevalence, intent proximity, severity, addressability and confidence.")

    tx = slide4.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(7.0), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "Intent × friction distribution (n = 1,500)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    dist_cards = [
        ("51.3%", "HIGH INTENT / HIGH FRICTION\nPrimary growth target (n=770)", ROSE, WHITE),
        ("22.0%", "LOW INTENT / HIGH FRICTION\nPassive browsing (n=330)", AMBER_BG, NAVY_DARK),
        ("16.7%", "HIGH INTENT / LOW FRICTION\nWaiting to execute (n=250)", GREEN_BG, NAVY_DARK),
        ("10.0%", "LOW INTENT / LOW FRICTION\nMood-board archive (n=150)", BG_OFFWHITE, NAVY_DARK),
    ]

    for idx, (num, label, bg_c, text_c) in enumerate(dist_cards):
        row = idx // 2; col = idx % 2
        c_left = 0.7 + col * 3.5; c_top = 2.4 + row * 1.9
        add_card(slide4, c_left, c_top, 3.3, 1.7, bg_color=bg_c, border_color=BORDER_COLOR if bg_c != ROSE else None)
        tx_box = slide4.shapes.add_textbox(Inches(c_left + 0.15), Inches(c_top + 0.15), Inches(3.0), Inches(1.4))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = text_c if bg_c != AMBER_BG and bg_c != GREEN_BG else (AMBER if bg_c == AMBER_BG else GREEN)
        p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = text_c

    tx = slide4.shapes.add_textbox(Inches(8.0), Inches(2.0), Inches(7.3), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "Opportunity ranking (/25)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    opp_headers = ["Opportunity Area", "Breakdown (P I S A C)", "Score & Rank"]
    opp_data = [
        [("1. Size & Fit Uncertainty", True, NAVY_DARK), "4  5  5  5  4", ("23 / 25 — #1", True, ROSE)],
        ["2. Comparison Paralysis", "5  5  4  5  3", "22 / 25 — #2"],
        ["3. Product Reality / Quality", "5  4  3  4  3", "19 / 25 — #3"],
        ["4. Price Volatility & Context", "4  4  4  2  3", "17 / 25 — #4"],
        ["5. Styling & Wardrobe Match", "3  4  3  3  3", "16 / 25 — #5"]
    ]
    add_table(slide4, 8.0, 2.4, 7.3, 3.6, opp_headers, opp_data, [3.6, 1.8, 1.9])

    add_card(slide4, 0.7, 6.4, 14.6, 1.5, bg_color=ROSE_BG, border_color=ROSE, left_accent_color=ROSE)
    tx_b = slide4.shapes.add_textbox(Inches(0.9), Inches(6.55), Inches(14.2), Inches(1.2))
    tf_b = tx_b.text_frame; tf_b.word_wrap = True
    p = tf_b.paragraphs[0]
    p.text = "Frequency ≠ priority. Fit wins because the affected users are closer to purchase, the friction is severe, and the problem is addressable without monetary incentives."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    add_footer(slide4, 4)

    # ==========================================
    # SLIDE 5: Primary Research
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "PRIMARY RESEARCH",
               "Interviews narrowed the hypothesis: fit appeared in 5/6 journeys and was the primary blocker for 3/6.",
               "Six wishlist walkthroughs reconstructed the path from “saved” to “researched” to “delayed / bought” for near-term purchase intent.")

    add_card(slide5, 0.7, 1.9, 4.6, 1.4, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=NAVY_DARK)
    tx = slide5.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(4.3), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "5 / 6"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY_DARK
    p = tf.add_paragraph(); p.text = "mentioned fit or drape in decision journey"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    add_card(slide5, 5.7, 1.9, 4.6, 1.4, bg_color=ROSE_BG, border_color=ROSE, left_accent_color=ROSE)
    tx = slide5.shapes.add_textbox(Inches(5.85), Inches(2.0), Inches(4.3), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "3 / 6"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = ROSE
    p = tf.add_paragraph(); p.text = "identified fit/size as primary blocker"; p.font.size = Pt(14); p.font.color.rgb = NAVY_DARK

    add_card(slide5, 10.7, 1.9, 4.6, 1.4, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=NAVY_DARK)
    tx = slide5.shapes.add_textbox(Inches(10.85), Inches(2.0), Inches(4.3), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "6 / 6"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY_DARK
    p = tf.add_paragraph(); p.text = "used an outside-app workaround"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    p_headers = ["P#", "Wishlisted item", "Primary blocker", "Purchase Trigger"]
    p_data = [
        ["P1", "Roadster jeans", ("Fit / size", True, ROSE), "Reliable fit comparison tool"],
        ["P2", "Anouk dress", ("Fit & drape", True, ROSE), "User photo on similar height model"],
        ["P3", "HRX shoes", "Comparison", "Side-by-side spec comparison"],
        ["P4", "Libas kurti set", ("Post-wash fit", True, ROSE), "Shrinkage rating guarantee"],
        ["P5", "Allen Solly blazer", "Fabric quality", "Unedited close-up fabric photo"],
        ["P6", "Highlander cargos", "Styling", "Virtual outfit pairing"]
    ]
    add_table(slide5, 0.7, 3.5, 6.8, 4.6, p_headers, p_data, [0.6, 1.8, 1.8, 2.6])

    add_card(slide5, 7.8, 3.5, 7.5, 4.6, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx_w = slide5.shapes.add_textbox(Inches(8.1), Inches(3.7), Inches(7.0), Inches(4.2))
    tf_w = tx_w.text_frame; tf_w.word_wrap = True
    p = tf_w.paragraphs[0]
    p.text = "Observed workaround loop:"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(10)

    loop_steps = [
        ("1. SAVE: ", "Specific item / size"),
        ("2. LEAVE APP: ", "Reddit · YouTube · Instagram"),
        ("3. TRIANGULATE: ", "Past orders · reviews · social proof"),
        ("4. DELAY: ", "Still not confident enough to commit")
    ]
    for ls_title, ls_desc in loop_steps:
        p = tf_w.add_paragraph()
        run1 = p.add_run(); run1.text = ls_title; run1.font.bold = True; run1.font.color.rgb = NAVY_DARK; run1.font.size = Pt(14)
        run2 = p.add_run(); run2.text = ls_desc; run2.font.color.rgb = SLATE_TEXT; run2.font.size = Pt(14)
        p.space_after = Pt(6)

    p = tf_w.add_paragraph()
    p.space_before = Pt(10)
    p.text = "“Roadster 30 is too tight; 32 slips off. I compared waist measurements from 3 previous orders.” — P1"
    p.font.size = Pt(14); p.font.bold = True; p.font.italic = True; p.font.color.rgb = ROSE

    add_footer(slide5, 5)

    # ==========================================
    # SLIDE 6: Problem Definition
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "PROBLEM DEFINITION",
               "The root problem is not lack of desire — it is lack of trusted, personalized fit confidence at the commitment moment.",
               "The thinking narrows consistently from the business metric to a specific user, moment, root cause and product outcome.")

    add_card(slide6, 0.7, 1.9, 14.6, 1.5, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER)
    tx_ps = slide6.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(14.2), Inches(1.2))
    tf_ps = tx_ps.text_frame; tf_ps.word_wrap = True
    p = tf_ps.paragraphs[0]
    run = p.add_run(); run.text = "Problem Statement: "; run.font.bold = True; run.font.color.rgb = AMBER; run.font.size = Pt(14)
    run = p.add_run()
    run.text = "High-intent apparel shoppers with a near-term purchase goal delay moving a wishlisted item to Bag because they cannot confidently predict how their selected size will fit across brands — so they leave Myntra to triangulate fit from reviews, past orders and try-on content before committing."
    run.font.color.rgb = WHITE; run.font.size = Pt(14)

    grid_cards_6 = [
        ("WHO", "High-intent apparel shoppers with a near-term purchase goal"),
        ("WHEN", "Revisiting a specific saved SKU to decide “buy now or wait?”"),
        ("ROOT CAUSE", "Size labels are not trusted as personal, cross-brand fit evidence"),
        ("WORKAROUND", "Reddit / YouTube / Instagram + comparing previous kept orders"),
        ("USER VALUE", "Decide fit without leaving Myntra or ordering “just to try”"),
        ("BUSINESS VALUE", "Increase Wishlist → Bag; fit accuracy may also reduce returns")
    ]

    for idx, (gc_title, gc_desc) in enumerate(grid_cards_6):
        row = idx // 3; col = idx % 3
        c_left = 0.7 + col * 4.96; c_top = 3.6 + row * 2.2
        add_card(slide6, c_left, c_top, 4.7, 2.1, bg_color=WHITE, border_color=BORDER_COLOR)
        tx_box = slide6.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.2), Inches(4.3), Inches(1.7))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = gc_title
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(6)
        p = tf.add_paragraph()
        p.text = gc_desc
        p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_footer(slide6, 6)

    # ==========================================
    # SLIDE 7: Solution Rationale
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "SOLUTION RATIONALE",
               "Smart Wishlist + FitCheck resolves the declared blocker inside the wishlist, then hands control back to the user.",
               "The core MVP is a decision-confidence loop — not a discount, coupon, or forced checkout nudge.")

    sol_steps = [
        ("1. FIT ANCHOR", "Confirm kept size: Roadster M", False),
        ("2. WAITING ON FIT", "Declare blocker: Waiting on Fit", False),
        ("3. FITCHECK", "Explain + recommend signal", False),
        ("4. READY TO BUY", "Resolve condition state", True),
        ("5. MOVE TO BAG", "Shopper commits & acts", False),
    ]
    for idx, (s_title, s_desc, is_hi) in enumerate(sol_steps):
        c_left = 0.7 + idx * (card_w + 0.2)
        bg_c = ROSE_BG if is_hi else WHITE
        brd_c = ROSE if is_hi else BORDER_COLOR
        add_card(slide7, c_left, 1.9, card_w, 2.0, bg_color=bg_c, border_color=brd_c)
        tx_box = slide7.shapes.add_textbox(Inches(c_left + 0.15), Inches(2.05), Inches(card_w - 0.3), Inches(1.7))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE if is_hi else NAVY_DARK
        p.space_after = Pt(6)
        p = tf.add_paragraph()
        p.text = s_desc
        p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide7, 0.7, 4.1, 7.1, 3.1, bg_color=GREEN_BG, border_color=GREEN)
    tx = slide7.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(6.7), Inches(2.7))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Why this fits the brief"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = GREEN; p.space_after = Pt(10)
    p = tf.add_paragraph()
    p.text = "No subsidy. No coupon. The value is informational: translate fit evidence into a personal recommendation, then let the shopper decide."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide7, 8.2, 4.1, 7.1, 3.1, bg_color=AMBER_CARD_BG, border_color=AMBER)
    tx = slide7.shapes.add_textbox(Inches(8.4), Inches(4.3), Inches(6.7), Inches(2.7))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Scope discipline"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(10)
    p = tf.add_paragraph()
    p.text = "Price / Both are expansion modes in the prototype. The causal experiment isolates Fit-only."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_button(slide7, 12.3, 7.4, 3.0, 0.65, "Launch live MVP 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    add_footer(slide7, 7)

    # ==========================================
    # SLIDE 8: Deployed MVP
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "DEPLOYED MVP",
               "The deployed MVP is intentionally thin: it proves the decision-resolution loop before investing in real fit ML.",
               "A top-scoring MVP is explicit about what is functional, what is simulated, and exactly what learning it unlocks.")

    add_card(slide8, 0.7, 1.9, 7.1, 4.9, bg_color=WHITE, border_color=BORDER_COLOR)
    tx = slide8.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(6.7), Inches(4.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Functional in the live prototype"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(10)

    func_items = [
        "✓ Set and save a Fit Anchor profile",
        "✓ Declare the blocker: Fit / Price / Both",
        "✓ Run a FitCheck decision event",
        "✓ Transition the item to “Ready to Buy”",
        "✓ Move the resolved item to Bag"
    ]
    for fi in func_items:
        p = tf.add_paragraph()
        p.text = fi
        p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT; p.space_after = Pt(6)

    add_button(slide8, 0.9, 6.0, 2.8, 0.65, "Test live MVP 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    add_card(slide8, 8.2, 1.9, 7.1, 4.9, bg_color=NAVY_CARD, border_color=None)
    tx = slide8.shapes.add_textbox(Inches(8.4), Inches(2.1), Inches(6.7), Inches(4.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Not production-ready yet — by design"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(10)

    non_func_items = [
        "• Real order / return-history ingestion",
        "• Production fit model + category calibration",
        "• Live inventory / price event integrations",
        "• Notification preferences + frequency caps",
        "• Learned confidence thresholds + abstention logic"
    ]
    for nfi in non_func_items:
        p = tf.add_paragraph()
        p.text = nfi
        p.font.size = Pt(14); p.font.color.rgb = BORDER_COLOR; p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    p.text = "Transparency prevents “prototype theatre”: MVP validates mechanism."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER

    add_card(slide8, 0.7, 7.0, 14.6, 1.0, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR)
    tx_b = slide8.shapes.add_textbox(Inches(0.9), Inches(7.15), Inches(14.2), Inches(0.7))
    tf_b = tx_b.text_frame
    p = tf_b.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "60-second evaluator path: [1] Set anchor → [2] Choose Fit → [3] Run FitCheck → [4] See Ready → [5] Move to Bag"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    add_footer(slide8, 8)

    # ==========================================
    # SLIDE 9: Define Success
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "DEFINE SUCCESS",
               "Success is a causal lift in 30-day wishlist buyers, with faster fit-to-bag progression and no return penalty.",
               "The metric stack separates business impact, mechanism validation, leading indicators and guardrails — each with a clear decision use.")

    top_metrics = [
        ("BUSINESS METRIC", "30-Day Wishlist Buyer Conversion", NAVY_CARD, WHITE, WHITE),
        ("PRIMARY MECHANISM", "Fit-Resolved → Bag within 72h", ROSE, WHITE, ROSE_BG),
        ("KEY GUARDRAIL", "Fit-related return rate", GREEN_BG, GREEN, NAVY_DARK)
    ]
    for idx, (m_tag, m_val, bg_c, tag_c, val_c) in enumerate(top_metrics):
        c_left = 0.7 + idx * 4.96
        add_card(slide9, c_left, 1.9, 4.7, 1.6, bg_color=bg_c, border_color=GREEN if bg_c == GREEN_BG else None)
        tx_box = slide9.shapes.add_textbox(Inches(c_left + 0.2), Inches(2.05), Inches(4.3), Inches(1.3))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_tag
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = tag_c; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = m_val
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = val_c

    m_headers = ["Metric", "Definition", "Why it matters"]
    m_data = [
        [("Anchor adoption", True, NAVY_DARK), "% fit-tagged users with a usable fit anchor", "Can users reach the value moment?"],
        [("Fit resolution rate", True, NAVY_DARK), "% FitChecks returning a decisive, explainable recommendation", "Does system actually resolve uncertainty?"],
        [("Median wishlist → Bag time", True, NAVY_DARK), "Hours from wishlist add / revisit to Bag for fit cohort", "Is decision latency shrinking?"],
        [("Wishlist add rate", True, NAVY_DARK), "Eligible users adding ≥1 item to wishlist", "Ensure extra friction does not suppress saving"],
        [("Recommendation disagreement", True, NAVY_DARK), "% users who override / reject the suggested size", "Early trust and calibration signal"]
    ]
    add_table(slide9, 0.7, 3.7, 14.6, 4.4, m_headers, m_data, [3.2, 6.0, 5.4])

    add_footer(slide9, 9)

    # ==========================================
    # SLIDE 10: Risks & Mitigation
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "RISKS & MITIGATION",
               "The biggest failure mode is false confidence; rollout should privilege abstention, trust and learning over reach.",
               "The solution fails if it pushes users faster into wrong-size orders, creates setup friction, or generalizes a narrow research sample too aggressively.")

    r_headers = ["Risk", "Impact", "Why it could happen", "Mitigation / rollout gate"]
    r_data = [
        [("False fit confidence → returns", True, NAVY_DARK), ("HIGH", True, ROSE), "Overconfident reasoning without measurement truth", "Confidence bands + \"not enough evidence\" state · pilot stable categories"],
        [("Fit Anchor setup friction", True, NAVY_DARK), ("MED", True, AMBER), "Manual multi-field setup delays value", "Auto-suggest from kept orders contextually"],
        [("Cold start / sparse brand data", True, NAVY_DARK), ("HIGH", True, ROSE), "New brands/categories lack evidence", "Blend size charts + material; abstain below threshold"],
        [("Research sample does not generalize", True, NAVY_DARK), ("MED", True, AMBER), "6 interviews are directional", "User-level A/B test before full rollout"],
        [("Notification fatigue / trust erosion", True, NAVY_DARK), ("MED", True, AMBER), "\"Ready to Buy\" nudges can feel pushy", "Trigger only on condition change · frequency cap"]
    ]
    add_table(slide10, 0.7, 1.9, 14.6, 4.5, r_headers, r_data, [3.6, 1.2, 4.8, 5.0])

    add_card(slide10, 0.7, 6.6, 14.6, 1.4, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER)
    tx_b10 = slide10.shapes.add_textbox(Inches(0.9), Inches(6.75), Inches(14.2), Inches(1.1))
    tf_b10 = tx_b10.text_frame; tf_b10.word_wrap = True
    p = tf_b10.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Rollout principle: Earn trust before reach: offline calibration → Fit-only pilot → category expansion. “A wishlist should manage the decision, not just remember the product.”"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE

    add_footer(slide10, 10)

    output_path = r"C:\Users\sanja\.gemini\antigravity\scratch\myntra_wishlist_conversion\Myntra_Wishlist_Growth_Case_Study.pptx"
    try:
        prs.save(output_path)
        print(f"Successfully generated PowerPoint presentation at: {output_path}")
    except Exception as e:
        print(f"Primary path failed: {e}")
        fallback_path = r"C:\Users\sanja\.gemini\antigravity\scratch\myntra_wishlist_conversion\Myntra_Wishlist_Growth_Case_Study_v2.pptx"
        prs.save(fallback_path)
        print(f"Successfully generated PowerPoint presentation at fallback: {fallback_path}")

if __name__ == "__main__":
    create_presentation()
