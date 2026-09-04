import sys
from pptx import Presentation

# Set UTF-8 stdout
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation('NL Myntra.pptx')
with open('extracted_edited_deck.txt', 'w', encoding='utf-8') as f:
    for idx, slide in enumerate(prs.slides, 1):
        f.write(f"\n=== SLIDE {idx} ===\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        f.write(f"  - {txt}\n")
            elif shape.has_table:
                table = shape.table
                f.write("  [TABLE]\n")
                for r_idx, row in enumerate(table.rows):
                    row_vals = [c.text.strip().replace('\n', ' ') for c in row.cells]
                    f.write(f"    Row {r_idx}: {' | '.join(row_vals)}\n")

print("Successfully extracted deck content to extracted_edited_deck.txt")
