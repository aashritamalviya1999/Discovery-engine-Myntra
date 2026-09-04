import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # Enterprise Board-Level Color Palette (High Contrast & Color-Blind Friendly)
    BG_OFFWHITE = RGBColor(248, 250, 252) # #f8fafc
    WHITE = RGBColor(255, 255, 255)
    NAVY_DARK = RGBColor(15, 23, 42) # #0f172a
    NAVY_CARD = RGBColor(27, 38, 79) # #1b264f
    ROSE = RGBColor(225, 29, 72) # #e11d48
    ROSE_BG = RGBColor(255, 241, 242) # #fff1f2
    AMBER = RGBColor(217, 119, 6) # #d97706
    AMBER_BG = RGBColor(254, 243, 199) # #fef3c7
    AMBER_CARD_BG = RGBColor(255, 251, 235) # #fffbeb
    GREEN = RGBColor(16, 185, 129) # #10b981
    GREEN_BG = RGBColor(236, 253, 245) # #ecfdf5
    SLATE_TEXT = RGBColor(51, 65, 85) # #334155
    MUTED_TEXT = RGBColor(71, 85, 105) # #475569
    LIGHT_MUTED = RGBColor(100, 116, 139) # #64748b
    BORDER_COLOR = RGBColor(203, 213, 225) # #cbd5e1

    FONT_FAMILY = "Inter"

    def add_bg(slide, color=BG_OFFWHITE):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, kicker_text, title_text, subtitle_text=None, top_start=0.4):
        tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(top_start), Inches(14.8), Inches(1.3))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Kicker
        p0 = tf.paragraphs[0]
        p0.text = kicker_text.upper()
        p0.font.name = FONT_FAMILY
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = ROSE
        p0.space_after = Pt(2)

        # Action Title
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(21)
        p1.font.bold = True
        p1.font.color.rgb = NAVY_DARK
        p1.space_after = Pt(2)

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(14)
            p2.font.color.rgb = MUTED_TEXT

    def add_footer(slide, slide_num):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(8.35), Inches(14.8), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER_COLOR
        line.line.fill.background()

        tx_left = slide.shapes.add_textbox(Inches(0.6), Inches(8.4), Inches(10), Inches(0.4))
        tf_left = tx_left.text_frame
        tf_left.margin_left = tf_left.margin_top = tf_left.margin_right = tf_left.margin_bottom = 0
        p_l = tf_left.paragraphs[0]
        p_l.text = "Myntra · Growth PM Case Study · Wishlist-to-Purchase Conversion"
        p_l.font.name = FONT_FAMILY
        p_l.font.size = Pt(14)
        p_l.font.color.rgb = MUTED_TEXT

        tx_right = slide.shapes.add_textbox(Inches(12.4), Inches(8.4), Inches(3.0), Inches(0.4))
        tf_right = tx_right.text_frame
        tf_right.margin_left = tf_right.margin_top = tf_right.margin_right = tf_right.margin_bottom = 0
        p_r = tf_right.paragraphs[0]
        p_r.text = f"Slide {slide_num}"
        p_r.alignment = PP_ALIGN.RIGHT
        p_r.font.name = FONT_FAMILY
        p_r.font.size = Pt(14)
        p_r.font.bold = True
        p_r.font.color.rgb = MUTED_TEXT

    def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=None, left_accent_width=0.08):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()

        if left_accent_color:
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(left_accent_width), Inches(height))
            accent.fill.solid()
            accent.fill.fore_color.rgb = left_accent_color
            accent.line.fill.background()
        return shape

    def add_button(slide, left, top, width, height, text, url, bg_color=NAVY_DARK):
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        btn.fill.solid()
        btn.fill.fore_color.rgb = bg_color
        btn.line.fill.background()

        tf = btn.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        if url:
            btn.click_action.hyperlink.address = url
        return btn

    def add_table(slide, left, top, width, height, headers, data, col_widths=None):
        rows = len(data) + 1
        cols = len(headers)
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = table_shape.table

        if col_widths:
            for idx, w in enumerate(col_widths):
                table.columns[idx].width = Inches(w)

        # Header Row
        for c_idx, h_text in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY_CARD
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = Inches(0.06)
            p = tf.paragraphs[0]
            p.text = h_text
            p.font.name = FONT_FAMILY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE

        # Data Rows
        for r_idx, row_data in enumerate(data):
            row_bg = WHITE if r_idx % 2 == 1 else BG_OFFWHITE
            for c_idx, val in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                tf = cell.text_frame
                tf.margin_left = tf.margin_right = Inches(0.12)
                tf.margin_top = tf.margin_bottom = Inches(0.05)
                p = tf.paragraphs[0]
                
                if isinstance(val, tuple):
                    p.text = str(val[0])
                    p.font.bold = val[1]
                    if len(val) > 2 and val[2]:
                        p.font.color.rgb = val[2]
                    else:
                        p.font.color.rgb = NAVY_DARK
                else:
                    p.text = str(val)
                    p.font.bold = False
                    p.font.color.rgb = NAVY_DARK

                p.font.name = FONT_FAMILY
                p.font.size = Pt(14)

        return table_shape

    # ==========================================
    # SLIDE 1: Cover & Executive Overview
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)
    add_header(slide1, "GROWTH PM CASE STUDY · MYNTRA",
               "Convert High-Intent Wishlisters by Resolving Size & Fit Uncertainty — Without Monetary Subsidies",
               "Wishlist adds are a stronger-than-browse interest signal, but purchase intent varies across user segments.")

    # Executive Summary Banner (Width 14.8in)
    add_card(slide1, 0.6, 1.7, 14.8, 1.3, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER, left_accent_width=0.1)
    tx_exec = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(14.4), Inches(1.1))
    tf_exec = tx_exec.text_frame; tf_exec.word_wrap = True
    p = tf_exec.paragraphs[0]
    run = p.add_run(); run.text = "EXECUTIVE SUMMARY & STRATEGIC THESIS: "; run.font.bold = True; run.font.color.rgb = AMBER; run.font.size = Pt(14)
    run = p.add_run()
    run.text = "Wishlist items represent high purchase desire, yet conversion stalls because shoppers cannot predict cross-brand fit drape without physically trying on garments. By deploying FitCheck — a personalized fit-confidence experience that uses a known-good past fit as an anchor to estimate confidence on a new item — Myntra bridges the gap from 'Saved' to 'Ready to Buy' without relying on discount subsidies or margin erosion."
    run.font.color.rgb = WHITE; run.font.size = Pt(14)

    # 4 Defensible Proof Point Cards (Width 3.55in each)
    stats = [
        ("51.3%", "High Intent / High Friction", "770 / 1,500 conversations fall into the segment closest to purchase but blocked by friction.", ROSE),
        ("23 / 25", "Opportunity Rubric #1", "Ranked #1 on intent proximity, severity, and non-monetary addressability.", AMBER),
        ("83% (5/6)", "Surfaced Fit/Drape Doubt", "5/6 interviewees surfaced fit uncertainty; 3/6 (50%) cited it as primary #1 blocker.", GREEN),
        ("6 / 6 (100%)", "Outside-App Workarounds", "All 6 participants used outside-app workarounds to resolve fit uncertainty before deciding.", NAVY_DARK),
    ]
    for idx, (s_num, s_title, s_desc, s_color) in enumerate(stats):
        c_left = 0.6 + idx * 3.75
        add_card(slide1, c_left, 3.2, 3.55, 2.0, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=s_color)
        tx = slide1.shapes.add_textbox(Inches(c_left + 0.15), Inches(3.3), Inches(3.25), Inches(1.8))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_num; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = s_color; p.space_after = Pt(2)
        p = tf.add_paragraph()
        p.text = s_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = s_desc; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    # Bottom Row: Process Flow & Artefact Action Links
    add_card(slide1, 0.6, 5.4, 6.8, 2.8, bg_color=WHITE, border_color=BORDER_COLOR)
    tx_l = slide1.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(6.4), Inches(2.5))
    tf_l = tx_l.text_frame; tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "TESTABLE ARTEFACTS & LIVE PROTOTYPES"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(8)
    p = tf_l.add_paragraph()
    p.text = "Publicly accessible tools built specifically for this case study submission:"; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT; p.space_after = Pt(12)

    add_button(slide1, 0.8, 6.6, 3.1, 0.65, "Open Discovery Engine 🔗", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/?v=10", bg_color=NAVY_DARK)
    add_button(slide1, 4.1, 6.6, 3.1, 0.65, "Open Live Deployed MVP 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    add_card(slide1, 7.6, 5.4, 7.8, 2.8, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=ROSE, left_accent_width=0.1)
    tx_r = slide1.shapes.add_textbox(Inches(7.85), Inches(5.55), Inches(7.3), Inches(2.5))
    tf_r = tx_r.text_frame; tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "END-TO-END CONVERSION FLOW"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(8)

    steps = [
        ("1. DISCOVERY", "1,500 public conversations classified into Intent × Friction quadrants."),
        ("2. RESEARCH", "6 qualitative walkthroughs isolate cross-brand size distrust."),
        ("3. FITCHECK ENGINE", "Garment attributes matched against user's known-good fit anchor."),
        ("4. READY TO BUY", "Wishlist item flips to action-ready state upon confidence signal."),
        ("5. 30D CONVERSION", "User moves item to Bag & signals stronger purchase commitment.")
    ]
    for st_title, st_desc in steps:
        p = tf_r.add_paragraph()
        run1 = p.add_run(); run1.text = f"{st_title}: "; run1.font.bold = True; run1.font.color.rgb = ROSE; run1.font.size = Pt(14)
        run2 = p.add_run(); run2.text = st_desc; run2.font.color.rgb = SLATE_TEXT; run2.font.size = Pt(14)
        p.space_after = Pt(2)

    add_footer(slide1, 1)

    # ==========================================
    # SLIDE 2: Business Metric Decomposition
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "BUSINESS METRIC DECOMPOSITION",
               "Wishlist-to-Bag Progression Is the Primary Growth Lever for 30-Day Conversion",
               "Decomposing the business North Star into controllable product outcomes, user drop-off stages, and non-goals.")

    # Navy Formula Banner (Matches Assignment Brief Exactly)
    add_card(slide2, 0.6, 1.7, 14.8, 1.0, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER, left_accent_width=0.08)
    tx_f = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(14.4), Inches(0.8))
    tf_f = tx_f.text_frame; tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "30-Day Wishlist Buyer Conversion = (Unique users purchasing ≥1 item they wishlisted within 30 days of adding it) ÷ (Unique cohort users with ≥1 eligible wishlist add)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE

    # Behavioral Drivers Bar
    add_card(slide2, 0.6, 2.8, 14.8, 0.5, bg_color=WHITE, border_color=BORDER_COLOR)
    tx_tree = slide2.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(14.4), Inches(0.4))
    tf_tree = tx_tree.text_frame; tf_tree.word_wrap = True
    p = tf_tree.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Behavioral drivers of 30D conversion: Wishlist Revisit → Intent Persistence → Decision Confidence → Wishlist to Bag → Bag to Purchase"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    # 5 Funnel Stage Cards (Baseline to Validate)
    funnel = [
        ("1. Wishlist Revisit", "Return to saved item ≤7 days", "Drop-off: ~25%*\nPassive retention", False),
        ("2. Decision Resolution", "Fit / drape uncertainty resolved", "Drop-off: ~48%*\nLARGEST BLOCKER", True),
        ("3. Wishlist → Bag", "Commit shortlisted SKU to Bag", "Target Lift: +14%*\nPrimary Outcome", True),
        ("4. Bag → Purchase", "Complete checkout transaction", "Drop-off: ~12%*\nStandard friction", False),
        ("5. 30-Day Cohort", "Complete purchase within window", "Cohort Lifetime\nIntent retention", False),
    ]
    card_w = 2.8
    for idx, (f_title, f_sub, f_drop, is_target) in enumerate(funnel):
        c_left = 0.6 + idx * (card_w + 0.2)
        bg_c = ROSE_BG if is_target else WHITE
        brd_c = ROSE if is_target else BORDER_COLOR
        add_card(slide2, c_left, 3.4, card_w, 1.8, bg_color=bg_c, border_color=brd_c)
        tx_box = slide2.shapes.add_textbox(Inches(c_left + 0.15), Inches(3.5), Inches(card_w - 0.3), Inches(1.6))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE if is_target else NAVY_DARK; p.space_after = Pt(2)
        p = tf.add_paragraph()
        p.text = f_sub; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = f_drop; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE if is_target else MUTED_TEXT

    # Bottom 2 Analytical Blocks (Width 7.3in each)
    add_card(slide2, 0.6, 5.3, 7.3, 2.9, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx = slide2.shapes.add_textbox(Inches(0.8), Inches(5.45), Inches(6.9), Inches(2.6))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STRATEGIC INTERVENTION FOCUS: WISHLIST → BAG"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "• For the fit-focused cohort, product interest is established; unresolved fit confidence remains the primary decision barrier.\n• Wishlist → Bag is the stage with the largest, most addressable drop-off (~48% of high-intent saved items are never added to Bag).\n• FitCheck targets the two adjacent, controllable friction steps: resolve fit doubt → accelerate Wishlist-to-Bag conversion without monetary discounts."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide2, 8.1, 5.3, 7.3, 2.9, bg_color=AMBER_CARD_BG, border_color=AMBER, left_accent_color=AMBER)
    tx = slide2.shapes.add_textbox(Inches(8.3), Inches(5.45), Inches(6.9), Inches(2.6))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EXPLICIT NON-GOALS & BOUNDARIES (*BASELINE TO VALIDATE)"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "• We do NOT optimize generic wishlist additions or total app time. More bookmarking is not business success unless intent converts within 30 days.\n• We do NOT issue price cuts, coupons, or cashback incentives. The strategic constraint requires solving user decision uncertainty natively.\n• *Note: Funnel percentages represent baseline hypotheses to validate after internal Myntra product analytics instrumentation."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_footer(slide2, 2)

    # ==========================================
    # SLIDE 3: AI-Powered Discovery Engine
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "AI-POWERED DISCOVERY ENGINE",
               "An Automated NLP Pipeline Classifying 1,500 Public Conversations into Intent × Friction Evidence",
               "Ingesting unstructured posts from Reddit, Play Store, and fashion forums into structured, decision-ready evidence objects.")

    # 5 Pipeline Cards (Explicit Rubric Dimensions P I S A C)
    proc_steps = [
        ("1. INGESTION", "1,500 public posts from Reddit, App/Play Store, YouTube reviews."),
        ("2. STRUCTURE", "14-field multi-label schema (Intent, Friction, Workaround, Need)."),
        ("3. SEGMENT", "LLM-tagged segmentation into 4 Intent × Friction quadrants."),
        ("4. PRIORITIZE", "Opportunity scoring: Prevalence + Intent + Severity + Addressability + Confidence /25."),
        ("5. AUDIT", "300 records manually audited; 90% AI–human label agreement."),
    ]
    for idx, (p_title, p_desc) in enumerate(proc_steps):
        c_left = 0.6 + idx * (card_w + 0.2)
        add_card(slide3, c_left, 1.8, card_w, 1.8, bg_color=WHITE, border_color=BORDER_COLOR)
        tx_box = slide3.shapes.add_textbox(Inches(c_left + 0.15), Inches(1.9), Inches(card_w - 0.3), Inches(1.6))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = p_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = p_desc; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    # Bottom 2 Cards: Structured Evidence Object & Quality Audit Box
    add_card(slide3, 0.6, 3.8, 8.8, 4.4, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx = slide3.shapes.add_textbox(Inches(0.8), Inches(3.95), Inches(8.4), Inches(4.1))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STRUCTURED EVIDENCE OBJECT SAMPLE (RECORD #842)"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(8)

    p = tf.add_paragraph(); p.space_after = Pt(10)
    run = p.add_run(); run.text = " Intent: HIGH "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = ROSE
    run = p.add_run(); run.text = "   Friction: SIZE / FIT "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = AMBER
    run = p.add_run(); run.text = "   Workaround: PAST ORDERS "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = GREEN
    run = p.add_run(); run.text = "   Source: REDDIT "; run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = NAVY_DARK

    p = tf.add_paragraph()
    p.text = "• Verbatim Post Snippet: \"I love these Roadster slim jeans on Myntra, but size 30 is too tight around waist and 32 slips down. I've left them saved for 2 weeks while comparing waist measurements against 3 previous orders.\"\n• Model Inference: Intent Proximity = 0.94 (Near-term purchase); Friction Severity = High; Solvability = Non-Monetary Fit Anchor.\n• Value to PM: Enables systematic quantitative comparison across friction categories while maintaining strict provenance back to raw user feedback."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide3, 9.6, 3.8, 5.8, 4.4, bg_color=NAVY_CARD, border_color=None)
    tx_q = slide3.shapes.add_textbox(Inches(9.8), Inches(3.95), Inches(5.4), Inches(4.1))
    tf_q = tx_q.text_frame; tf_q.word_wrap = True
    p = tf_q.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "QUALITY CONTROL & MODEL AUDIT"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(8)

    p = tf_q.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    p.text = "300 Records Manually Audited"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.space_after = Pt(6)

    p = tf_q.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    p.text = "270 / 300 Matched = 90% Agreement"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BG_OFFWHITE; p.space_after = Pt(12)

    p = tf_q.add_paragraph()
    p.text = "• Measured: 270/300 records matched on primary Intent + Friction labels (90% AI–human label agreement).\n• Stack: Python + GPT-4o API + GitHub Pages Explorer\n• Provenance links maintained for full auditability."
    p.font.size = Pt(14); p.font.color.rgb = BORDER_COLOR; p.space_after = Pt(16)

    add_button(slide3, 10.7, 6.7, 3.6, 0.65, "Test Discovery Engine Live 🔗", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/?v=10", bg_color=ROSE)

    add_footer(slide3, 3)

    # ==========================================
    # SLIDE 4: Discovery Engine Findings
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "DISCOVERY ENGINE FINDINGS",
               "Size & Fit Uncertainty Ranks #1 — High Purchase Proximity and Non-Monetary Addressability",
               "Segmenting 1,500 conversations by Intent × Friction and scoring opportunity areas on a 25-point PM rubric.")

    # Left Column: 2x2 Matrix Cards
    tx = slide4.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(7.0), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "INTENT × FRICTION SEGMENTATION (n = 1,500)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    dist_cards = [
        ("51.3%", "HIGH INTENT / HIGH FRICTION\n51.3% (n=770) fall into High Intent / High Friction — the primary segment for conversion intervention.", ROSE, WHITE),
        ("22.0%", "LOW INTENT / HIGH FRICTION\nPassive Browsing (n=330): Casual window shopping & price watching.", AMBER_BG, NAVY_DARK),
        ("16.7%", "HIGH INTENT / LOW FRICTION\nWaiting to Execute (n=250): High desire, awaiting payday/event date.", GREEN_BG, NAVY_DARK),
        ("10.0%", "LOW INTENT / LOW FRICTION\nMood-Board Archive (n=150): Long-term outfit bookmarking.", BG_OFFWHITE, NAVY_DARK),
    ]

    for idx, (num, label, bg_c, text_c) in enumerate(dist_cards):
        row = idx // 2; col = idx % 2
        c_left = 0.6 + col * 3.5; c_top = 2.2 + row * 2.0
        add_card(slide4, c_left, c_top, 3.3, 1.8, bg_color=bg_c, border_color=BORDER_COLOR if bg_c != ROSE else None)
        tx_box = slide4.shapes.add_textbox(Inches(c_left + 0.15), Inches(c_top + 0.15), Inches(3.0), Inches(1.5))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = text_c if bg_c != AMBER_BG and bg_c != GREEN_BG else (AMBER if bg_c == AMBER_BG else GREEN)
        p.space_after = Pt(2)
        p = tf.add_paragraph()
        p.text = label; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = text_c

    # Right Column: Opportunity Ranking Table
    tx = slide4.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(7.6), Inches(0.4))
    p = tx.text_frame.paragraphs[0]
    p.text = "OPPORTUNITY PRIORITIZATION RUBRIC (P · I · S · A · C /25 TOTAL)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    opp_headers = ["Opportunity Area", "P  I  S  A  C Breakdown", "Score & Status"]
    opp_data = [
        [("1. Size & Fit Uncertainty", True, NAVY_DARK), "P:4  I:5  S:5  A:5  C:4", ("23 / 25 — #1 (PRIMARY WEDGE)", True, ROSE)],
        ["2. Comparison Paralysis", "P:5  I:5  S:4  A:5  C:3", "22 / 25 — #2 (Backlog)"],
        ["3. Product Quality / Real Photo Void", "P:5  I:4  S:3  A:4  C:3", "19 / 25 — #3 (Backlog)"],
        ["4. Price Volatility & Context", "P:4  I:4  S:4  A:2  C:3", "17 / 25 — #4 (Monetary Boundary)"],
        ["5. Styling & Wardrobe Match", "P:3  I:4  S:3  A:3  C:3", "16 / 25 — #5 (Low Priority)"]
    ]
    add_table(slide4, 7.8, 2.2, 7.6, 3.6, opp_headers, opp_data, [3.6, 2.0, 2.0])

    # Rubric Legend Bar
    add_card(slide4, 7.8, 5.9, 7.6, 0.4, bg_color=WHITE, border_color=BORDER_COLOR)
    tx_leg = slide4.shapes.add_textbox(Inches(7.9), Inches(5.95), Inches(7.4), Inches(0.3))
    tf_leg = tx_leg.text_frame
    p = tf_leg.paragraphs[0]
    p.text = "Legend: P = Prevalence · I = Intent proximity · S = Severity · A = Addressability · C = Confidence (1–5 each)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = SLATE_TEXT

    # Bottom Banner
    add_card(slide4, 0.6, 6.4, 14.8, 1.5, bg_color=ROSE_BG, border_color=ROSE, left_accent_color=ROSE)
    tx_b = slide4.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(14.4), Inches(1.2))
    tf_b = tx_b.text_frame; tf_b.word_wrap = True
    p = tf_b.paragraphs[0]
    p.text = "KEY TAKEAWAY: FREQUENCY ≠ PRIORITY FOR GROWTH INTERVENTION\nComparison friction appears slightly more often across browsing, but Size & Fit Uncertainty ranks #1 overall because affected users are closest to purchase, their friction is severe, and the problem is addressable without monetary discounts."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    add_footer(slide4, 4)

    # ==========================================
    # SLIDE 5: Primary Research
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "PRIMARY RESEARCH",
               "6 Wishlist Walkthroughs Reveal Users Leave Myntra to Resolve Fit Uncertainty",
               "Reconstructing user decision paths from 'Saved' to 'Researched' to 'Stalled / Purchased' for near-term purchase intent.")

    # 3 Defensible Stat Cards
    add_card(slide5, 0.6, 1.8, 4.7, 1.3, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=NAVY_DARK)
    tx = slide5.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(4.4), Inches(1.1))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "5 / 6 (83%)"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY_DARK
    p = tf.add_paragraph(); p.text = "surfaced size/drape uncertainty in decision journey"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    add_card(slide5, 5.65, 1.8, 4.7, 1.3, bg_color=ROSE_BG, border_color=ROSE, left_accent_color=ROSE)
    tx = slide5.shapes.add_textbox(Inches(5.8), Inches(1.85), Inches(4.4), Inches(1.1))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "3 / 6 (50%)"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = ROSE
    p = tf.add_paragraph(); p.text = "named cross-brand fit as their primary blocker"; p.font.size = Pt(14); p.font.color.rgb = NAVY_DARK

    add_card(slide5, 10.7, 1.8, 4.7, 1.3, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR, left_accent_color=NAVY_DARK)
    tx = slide5.shapes.add_textbox(Inches(10.85), Inches(1.85), Inches(4.4), Inches(1.1))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.text = "6 / 6 (100%)"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = NAVY_DARK
    p = tf.add_paragraph(); p.text = "used outside-app workarounds to resolve uncertainty"; p.font.size = Pt(14); p.font.color.rgb = MUTED_TEXT

    # User Decision Journey Flowchart Bar
    add_card(slide5, 0.6, 3.2, 14.8, 0.45, bg_color=WHITE, border_color=BORDER_COLOR)
    tx_j = slide5.shapes.add_textbox(Inches(0.8), Inches(3.23), Inches(14.4), Inches(0.4))
    tf_j = tx_j.text_frame; tf_j.word_wrap = True
    p = tf_j.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "USER DECISION JOURNEY: Wishlist Add → 'Will this fit me?' → External Research → Still Uncertain → Delayed Purchase"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    # Bottom Grid: Compact User Interview Matrix Table & Observed Loop Box (Zero Footer Overflow!)
    p_headers = ["P#", "Profile", "Saved SKU", "Stated Blocker", "Outside Workaround", "Required Purchase Trigger"]
    p_data = [
        ["P1", "Male (24)", "Roadster Jeans", ("Fit / size (30 vs 32)", True, ROSE), "Compared 3 past orders", "Reliable fit comparison tool"],
        ["P2", "Female (26)", "Anouk Kurta Set", ("Bust/shoulder drape", True, ROSE), "YouTube try-on hauls", "Real model height/body photos"],
        ["P3", "Male (22)", "HRX Shoes", "Foot width / sizing", "Reddit sizing threads", "Side-by-side spec comparison"],
        ["P4", "Female (29)", "Libas Kurti Set", ("Post-wash shrinkage", True, ROSE), "Read negative reviews", "Shrinkage rating guarantee"],
        ["P5", "Male (27)", "Allen Solly Blazer", "Arm length / shoulder", "Ordered 2 sizes to return 1", "Precise shoulder anchor"],
        ["P6", "Male (23)", "Highlander Cargos", "Waist fit / length", "Asked friends on WhatsApp", "Virtual outfit drape preview"]
    ]
    add_table(slide5, 0.6, 3.75, 8.8, 4.2, p_headers, p_data, [0.5, 1.1, 1.6, 1.8, 1.8, 2.0])

    add_card(slide5, 9.6, 3.75, 5.8, 4.2, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx_w = slide5.shapes.add_textbox(Inches(9.8), Inches(3.85), Inches(5.4), Inches(3.9))
    tf_w = tx_w.text_frame; tf_w.word_wrap = True
    p = tf_w.paragraphs[0]
    p.text = "USERS LEAVE MYNTRA TO RESOLVE FIT UNCERTAINTY"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(4)

    loop_steps = [
        ("1. SAVE ITEM: ", "User saves specific SKU & size in wishlist."),
        ("2. EXIT MYNTRA: ", "Leaves app to search Reddit, YouTube, Instagram."),
        ("3. TRIANGULATE: ", "Compares past kept orders & customer reviews."),
        ("4. STALL CHECKOUT: ", "Remains uncertain, delays purchase, intent decays.")
    ]
    for ls_title, ls_desc in loop_steps:
        p = tf_w.add_paragraph()
        run1 = p.add_run(); run1.text = ls_title; run1.font.bold = True; run1.font.color.rgb = NAVY_DARK; run1.font.size = Pt(14)
        run2 = p.add_run(); run2.text = ls_desc; run2.font.color.rgb = SLATE_TEXT; run2.font.size = Pt(14)
        p.space_after = Pt(2)

    p = tf_w.add_paragraph()
    p.space_before = Pt(4)
    p.text = "“I love this pair of Roadster jeans, but my waist size is 31 — Roadster 30 is too tight and 32 slips off. I left it saved for 2 weeks while comparing waist measurements from 3 previous kept orders.” — P1 (Walkthrough Interview)"
    p.font.size = Pt(14); p.font.bold = True; p.font.italic = True; p.font.color.rgb = ROSE

    add_footer(slide5, 5)

    # ==========================================
    # SLIDE 6: Problem Definition
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "PROBLEM DEFINITION",
               "High-Intent Shoppers Stall at Commitment Because Cross-Brand Size Labels Don't Translate into Personal Fit Confidence",
               "Synthesizing findings across Business Metric → Product Outcome → AI Discovery → Primary Research → Problem Statement.")

    # Top Evidence Chain Visual Strip
    add_card(slide6, 0.6, 1.7, 14.8, 0.5, bg_color=WHITE, border_color=BORDER_COLOR)
    tx_chain = slide6.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(14.4), Inches(0.4))
    tf_chain = tx_chain.text_frame; tf_chain.word_wrap = True
    p = tf_chain.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "EVIDENCE CHAIN: 30D Buyer Conversion → Wishlist to Bag → AI Discovery (Fit #1 23/25) → Interviews (5/6 Surface Fit) → Root Cause (Cross-Brand Size Distrust)"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    # Main Problem Statement Banner (Navy Dark, width 14.8in)
    add_card(slide6, 0.6, 2.3, 14.8, 1.3, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER, left_accent_width=0.1)
    tx_ps = slide6.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(14.4), Inches(1.1))
    tf_ps = tx_ps.text_frame; tf_ps.word_wrap = True
    p = tf_ps.paragraphs[0]
    run = p.add_run(); run.text = "CORE PROBLEM STATEMENT: "; run.font.bold = True; run.font.color.rgb = AMBER; run.font.size = Pt(14)
    run = p.add_run()
    run.text = "High-intent apparel shoppers with a near-term purchase goal delay moving wishlisted items to Bag because they cannot confidently predict how their selected size will fit across different brands — forcing them to leave Myntra to triangulate fit from external reviews, past purchases, and try-on content before committing."
    run.font.color.rgb = WHITE; run.font.size = Pt(14)

    # 6 Detailed Dimension Cards (3x2 Grid, Width 4.76in each)
    grid_cards_6 = [
        ("TARGET USER SEGMENT", "High-intent Myntra apparel shoppers who have shortlisted specific SKUs with near-term purchase intent."),
        ("DECISION TRIGGER", "Revisiting saved items to evaluate checkout readiness for an upcoming event, season change, or wardrobe update."),
        ("CORE ROOT OBSTACLE", "Inability to translate standard size labels (S/M/L or 30/32) into personalized fit guidance."),
        ("EXISTING WORKAROUND", "Leaving Myntra to browse Reddit sizing threads, Instagram try-ons, YouTube hauls, or measuring past orders."),
        ("EXPECTED USER VALUE", "Confidently resolve fit & size decisions quickly in one in-app flow without ordering multiple sizes to return."),
        ("EXPECTED BUSINESS VALUE", "Increases Wishlist-to-Bag conversion rate, shortens decision latency, and reduces fit-related returns downstream.")
    ]

    for idx, (gc_title, gc_desc) in enumerate(grid_cards_6):
        row = idx // 3; col = idx % 3
        c_left = 0.6 + col * 5.02; c_top = 3.7 + row * 2.2
        add_card(slide6, c_left, c_top, 4.76, 2.0, bg_color=WHITE, border_color=BORDER_COLOR)
        tx_box = slide6.shapes.add_textbox(Inches(c_left + 0.2), Inches(c_top + 0.2), Inches(4.36), Inches(1.6))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = gc_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = gc_desc; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_footer(slide6, 6)

    # ==========================================
    # SLIDE 7: Solution Rationale
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "SOLUTION RATIONALE",
               "Smart Wishlist + FitCheck Resolves Declared Blocker Inside the Wishlist — No Discounts Required",
               "Designing a native decision-confidence loop that empowers shoppers to declare obstacles and receive personalized fit recommendations.")

    # 5 Process Step Cards
    sol_steps = [
        ("1. FIT ANCHOR", "User confirms 1 kept item profile (e.g., Roadster M)."),
        ("2. DECLARE BLOCKER", "User tags SKU: 'Waiting on Fit'."),
        ("3. FITCHECK ENGINE", "FitCheck compares available product attributes with user's known-good fit anchor."),
        ("4. READY TO BUY", "Wishlist state transitions upon confidence signal."),
        ("5. MOVE TO BAG", "User moves item to Bag, signaling stronger purchase commitment."),
    ]
    for idx, (s_title, s_desc, is_target) in enumerate([(s[0], s[1], idx == 3) for idx, s in enumerate(sol_steps)]):
        c_left = 0.6 + idx * (card_w + 0.2)
        bg_c = ROSE_BG if is_target else WHITE
        brd_c = ROSE if is_target else BORDER_COLOR
        add_card(slide7, c_left, 1.8, card_w, 2.0, bg_color=bg_c, border_color=brd_c)
        tx_box = slide7.shapes.add_textbox(Inches(c_left + 0.15), Inches(1.9), Inches(card_w - 0.3), Inches(1.7))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE if is_target else NAVY_DARK; p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = s_desc; p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    # Bottom 2 Detailed Analysis Blocks
    add_card(slide7, 0.6, 4.0, 7.3, 3.2, bg_color=GREEN_BG, border_color=GREEN, left_accent_color=GREEN)
    tx = slide7.shapes.add_textbox(Inches(0.8), Inches(4.15), Inches(6.9), Inches(2.9))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "WHY THIS FITS THE NON-MONETARY STRATEGIC BRIEF"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = GREEN; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "• Zero Subsidy & Zero Margin Erosion: The solution does not rely on price slashes, coupons, or cashback nudges.\n• Pure Informational Value: Converts raw product reviews and kept order data into personal size confidence, solving the exact friction that stalls high-intent shoppers.\n• Handing Control to User: Users explicitly state what they are waiting for, transforming wishlist from passive storage into an active decision dashboard."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_card(slide7, 8.1, 4.0, 7.3, 3.2, bg_color=AMBER_CARD_BG, border_color=AMBER, left_accent_color=AMBER)
    tx = slide7.shapes.add_textbox(Inches(8.3), Inches(4.15), Inches(6.9), Inches(2.9))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SCOPE DISCIPLINE & EXPERIMENTAL CONTROL"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(8)
    p = tf.add_paragraph()
    p.text = "• Isolated Causal Experiment: The FitCheck MVP isolates Fit-only to measure pure decision confidence uplift.\n• Architecture Extensibility: Architecture can later support other declared blockers such as organic price events; excluded from the Fit MVP experiment.\n• Testable Artefact: Prototype deployed and accessible live."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    add_button(slide7, 12.0, 7.4, 3.4, 0.65, "Launch Live MVP Prototype 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    add_footer(slide7, 7)

    # ==========================================
    # SLIDE 8: Deployed MVP Prototype
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "DEPLOYED MVP PROTOTYPE",
               "Demonstrating the Decision-Resolution & State-Transition Mechanism",
               "Testing the non-monetary decision-confidence loop before investing engineering in real-time fit ML models.")

    # Left Column: Visual UI Wireframe Demonstrator Representation (Width 7.3in)
    add_card(slide8, 0.6, 1.8, 7.3, 5.2, bg_color=WHITE, border_color=BORDER_COLOR, left_accent_color=ROSE)
    tx_ui = slide8.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(6.9), Inches(4.8))
    tf_ui = tx_ui.text_frame; tf_ui.word_wrap = True
    p = tf_ui.paragraphs[0]
    p.text = "INTERACTIVE MVP UI EXPERIENCE & FLOW"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY_DARK; p.space_after = Pt(10)

    # UI Step 1 Box
    p = tf_ui.add_paragraph()
    run = p.add_run(); run.text = "1. FIT ANCHOR: "; run.font.bold = True; run.font.color.rgb = ROSE; run.font.size = Pt(14)
    run = p.add_run(); run.text = "User inputs known-good fit item: [Roadster Slim Jeans - Size 30 M]. Saved as personal reference anchor."; run.font.color.rgb = SLATE_TEXT; run.font.size = Pt(14)
    p.space_after = Pt(8)

    # UI Step 2 Box
    p = tf_ui.add_paragraph()
    run = p.add_run(); run.text = "2. FITCHECK SIGNAL: "; run.font.bold = True; run.font.color.rgb = AMBER; run.font.size = Pt(14)
    run = p.add_run(); run.text = "User taps 'Run FitCheck' on saved Anouk Dress. System compares garment specs against anchor → Returns 94% High Confidence match."; run.font.color.rgb = SLATE_TEXT; run.font.size = Pt(14)
    p.space_after = Pt(8)

    # UI Step 3 Box
    p = tf_ui.add_paragraph()
    run = p.add_run(); run.text = "3. READY TO BUY STATE: "; run.font.bold = True; run.font.color.rgb = GREEN; run.font.size = Pt(14)
    run = p.add_run(); run.text = "Item status flips from 'Waiting on Fit' to 'Ready to Buy'. One-tap 'Move to Bag' CTA unlocked without monetary discounts."; run.font.color.rgb = SLATE_TEXT; run.font.size = Pt(14)
    p.space_after = Pt(12)

    add_button(slide8, 0.8, 6.2, 3.2, 0.65, "Test Live Deployed MVP 🚀", "https://aashritamalviya1999.github.io/Discovery-engine-Myntra/fitcheck.html?v=10", bg_color=ROSE)

    # Right Column: Live Today vs Simulated Scope
    add_card(slide8, 8.1, 1.8, 7.3, 5.2, bg_color=NAVY_CARD, border_color=None)
    tx_sc = slide8.shapes.add_textbox(Inches(8.3), Inches(1.95), Inches(6.9), Inches(4.8))
    tf_sc = tx_sc.text_frame; tf_sc.word_wrap = True
    p = tf_sc.paragraphs[0]
    p.text = "LIVE TODAY IN MVP vs SIMULATED FOR PRODUCTION"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(10)

    items_live = [
        "LIVE TODAY IN MVP:",
        "✓ Set & Save Fit Anchor Profile (Brand, Size, Height, Body preference)",
        "✓ Declare Blocker Tag on Wishlist SKUs (Waiting on Fit / Price / Both)",
        "✓ Run FitCheck and generate a fit-confidence result",
        "✓ Transition Wishlist Item State from 'Waiting on Fit' to 'Ready to Buy'",
        "✓ Execute One-Tap 'Move to Bag' Action for Resolved Items",
        "",
        "SIMULATED / PRODUCTION BACKLOG:",
        "• Real-time kept order & return history API ingestion",
        "• Multi-category production ML garment measurement calibration model"
    ]
    for it in items_live:
        p = tf_sc.add_paragraph()
        if it.startswith("LIVE TODAY") or it.startswith("SIMULATED"):
            p.text = it; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER; p.space_after = Pt(4)
        else:
            p.text = it; p.font.size = Pt(14); p.font.color.rgb = BORDER_COLOR; p.space_after = Pt(4)

    p = tf_sc.add_paragraph(); p.space_before = Pt(6)
    p.text = "MVP validates the core decision loop and state transition; production ML accuracy remains an explicit validation step."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = AMBER

    # Bottom 60-Second Evaluator Path Banner
    add_card(slide8, 0.6, 7.1, 14.8, 1.0, bg_color=BG_OFFWHITE, border_color=BORDER_COLOR)
    tx_b = slide8.shapes.add_textbox(Inches(0.8), Inches(7.25), Inches(14.4), Inches(0.7))
    tf_b = tx_b.text_frame
    p = tf_b.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "60-SECOND EVALUATOR TESTING PATH: [1] Set Fit Anchor → [2] Select Item & Tag 'Waiting on Fit' → [3] Click FitCheck → [4] See 'Ready to Buy' State → [5] Move to Bag"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY_DARK

    add_footer(slide8, 8)

    # ==========================================
    # SLIDE 9: Define Success & A/B Experiment
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "DEFINE SUCCESS",
               "We scale only if FitCheck lifts 30-day conversion without increasing fit-related returns.",
               "Structuring metrics into business impact, mechanism validation, leading indicators, and guardrails.")

    # 3 Top Metric Stack Cards (Width 4.76in each)
    top_metrics = [
        ("BUSINESS NORTH STAR METRIC", "30-Day Wishlist Buyer Conversion Rate (% unique wishlist users purchasing ≥1 saved SKU)", NAVY_CARD, WHITE, WHITE),
        ("PRIMARY MECHANISM METRIC", "Condition-Resolved to Bag Rate (% of fit-tagged items moved to Bag within 72h)", ROSE, WHITE, ROSE_BG),
        ("KEY GUARDRAIL METRIC", "Fit-Related Return Rate (Must remain flat or drop; ensures fit accuracy)", GREEN_BG, GREEN, NAVY_DARK)
    ]
    for idx, (m_tag, m_val, bg_c, tag_c, val_c) in enumerate(top_metrics):
        c_left = 0.6 + idx * 5.02
        add_card(slide9, c_left, 1.8, 4.76, 1.5, bg_color=bg_c, border_color=GREEN if bg_c == GREEN_BG else None)
        tx_box = slide9.shapes.add_textbox(Inches(c_left + 0.2), Inches(1.85), Inches(4.36), Inches(1.3))
        tf = tx_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_tag; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = tag_c; p.space_after = Pt(2)
        p = tf.add_paragraph()
        p.text = m_val; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = val_c

    # A/B Test Experimentation Box (Width 14.8in)
    add_card(slide9, 0.6, 3.4, 14.8, 1.2, bg_color=WHITE, border_color=ROSE, left_accent_color=ROSE)
    tx_exp = slide9.shapes.add_textbox(Inches(0.8), Inches(3.45), Inches(14.4), Inches(1.1))
    tf_exp = tx_exp.text_frame; tf_exp.word_wrap = True
    p = tf_exp.paragraphs[0]
    p.text = "CAUSAL A/B TEST EXPERIMENT DESIGN & SAMPLE RIGOR"; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE; p.space_after = Pt(4)
    p = tf_exp.add_paragraph()
    p.text = "• Population: High-intent apparel wishlisters with ≥1 eligible fit-sensitive SKU   • Control: Existing Wishlist   • Treatment: Wishlist + FitCheck\n• Randomization: User-level A/B split   • Evaluation Window: 30 days   • Sample Size: Determined through power analysis using Myntra's baseline conversion and predefined minimum detectable effect (MDE).\n• Decision Rule: Launch only if 30-Day Conversion rises without inflating fit returns."
    p.font.size = Pt(14); p.font.color.rgb = SLATE_TEXT

    # Metrics Taxonomy Table (Corrected Metric Names & Rationales!)
    m_headers = ["Category", "Metric Name", "Exact Definition / Formula", "Target / PM Rationale"]
    m_data = [
        ["Leading Indicator", ("Fit Anchor Adoption Rate", True, NAVY_DARK), "% fit-tagged users with a saved anchor profile", "Pilot hypothesis — to be calibrated post-baseline"],
        ["Leading Indicator", ("Blocker Declaration Rate", True, NAVY_DARK), "% saved items tagged with specific blocker", "Pilot hypothesis — to be calibrated post-baseline"],
        ["Leading Indicator", ("Fit Resolution Coverage", True, NAVY_DARK), "% eligible FitChecks receiving recommendations", "Measures system ability to return confident calls"],
        ["Leading Indicator", ("Post-Purchase Fit Agreement", True, NAVY_DARK), "% purchases where users report fit was accurate", "Direct measure of fit recommendation precision"],
        ["Guardrail Metric", ("Wishlist Add Suppression", True, NAVY_DARK), "Total wishlist additions per user cohort", "Must not decrease (verifies no added friction)"],
        ["Guardrail Metric", ("Notification Opt-Out Rate", True, NAVY_DARK), "% users disabling FitCheck reminders", "Prevents notification fatigue and trust erosion"]
    ]
    add_table(slide9, 0.6, 4.7, 14.8, 3.5, m_headers, m_data, [1.8, 3.2, 5.4, 4.4])

    add_footer(slide9, 9)

    # ==========================================
    # SLIDE 10: Risk Analysis & Mitigation
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "RISKS & MITIGATION",
               "Mitigating False Fit Confidence, Setup Friction, and Cold Start Through Phased Rollout",
               "Evaluating potential failure modes, business risks, and defining strict rollout gating principles.")

    # Top Highlight Risk Banner
    add_card(slide10, 0.6, 1.7, 14.8, 0.6, bg_color=ROSE_BG, border_color=ROSE, left_accent_color=ROSE)
    tx_hr = slide10.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(14.4), Inches(0.5))
    tf_hr = tx_hr.text_frame; tf_hr.word_wrap = True
    p = tf_hr.paragraphs[0]
    p.text = "CRITICAL RISK: False fit confidence can turn a stalled wishlist into a costly return — worse than doing nothing."
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ROSE

    # Detailed Risk Matrix Table (5 Rows)
    r_headers = ["Risk Area", "Impact", "Likelihood", "Root Cause", "Mitigation Strategy & Rollout Gate"]
    r_data = [
        [("False Fit Confidence → High Returns", True, NAVY_DARK), ("HIGH", True, ROSE), "MED", "AI overconfident without garment specs", "Abstain below a calibrated confidence threshold; show confidence bands; gate rollout on return rate."],
        [("Fit Anchor Setup Friction", True, NAVY_DARK), "MED", ("HIGH", True, AMBER), "Manual multi-field input deters users", "Auto-infer anchor from kept orders; ask contextually when user adds item to wishlist."],
        [("Cold Start / Sparse Brand Data", True, NAVY_DARK), ("HIGH", True, ROSE), "MED", "New brands or niche categories lack fit evidence", "Blend size charts + material elasticity; output 'Not Enough Evidence' state instead of guessing."],
        [("Research Sample Generalization", True, NAVY_DARK), "MED", "MED", "6 qualitative interviews are directional", "Determine sample size through power analysis using baseline conversion and predefined MDE; segment by category."],
        [("Notification Fatigue / Nudge Erosion", True, NAVY_DARK), "MED", "MED", "Repeated 'Ready to Buy' alerts feel spammy", "Cap reminders to 1 per week per user; digest updates into weekly wishlist recap; 1-tap opt-out."]
    ]
    add_table(slide10, 0.6, 2.4, 14.8, 4.2, r_headers, r_data, [3.2, 1.0, 1.1, 4.0, 5.5])

    # Bottom Navy Rollout Principle Banner
    add_card(slide10, 0.6, 6.7, 14.8, 1.4, bg_color=NAVY_CARD, border_color=None, left_accent_color=AMBER)
    tx_b10 = slide10.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(14.4), Inches(1.1))
    tf_b10 = tx_b10.text_frame; tf_b10.word_wrap = True
    p = tf_b10.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "ROLLOUT PRINCIPLE: EARN TRUST BEFORE REACH\nOffline data calibration → Pilot in stable categories (Jeans & Ethnicwear) → Validate return guardrails → Category expansion. 'A wishlist should actively resolve decisions, not just store products.'"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE

    add_footer(slide10, 10)

    output_pptx = r"C:\Users\sanja\.gemini\antigravity\scratch\myntra_wishlist_conversion\NL_Myntra_TopScholar.pptx"
    try:
        prs.save(output_pptx)
        print(f"Successfully generated PowerPoint presentation at: {output_pptx}")
    except Exception as e:
        print(f"Primary save failed: {e}")

    fallback_pptx = r"C:\Users\sanja\.gemini\antigravity\scratch\myntra_wishlist_conversion\NL Myntra.pptx"
    try:
        prs.save(fallback_pptx)
        print(f"Successfully saved to fallback: {fallback_pptx}")
    except Exception as e:
        print(f"Fallback save failed: {e}")

if __name__ == "__main__":
    create_presentation()
